"""
/api/generate-pptx.py
Vercel Python serverless function.
Receives Agent 5 finalSpec JSON → builds PPTX using python-pptx → returns base64.

Input (POST body JSON):
  {
    "finalSpec": [...],       # Agent 5 designed spec array
    "brandRulebook": {...}    # Agent 2 brand rulebook (for fallback values)
  }

Output (JSON):
  {
    "success": true,
    "data": "<base64 pptx>",
    "slides": 12,
    "filename": "presentation_20241115.pptx"
  }
"""

import json
import base64
import io
import traceback
from datetime import datetime
from http.server import BaseHTTPRequestHandler

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree


# Renderer-side spacing / fit rules. These must not override Agent 5 placement,
# only how content is laid out inside the frames Agent 5 already chose.
INTERNAL_PADDING = 0.14
MIN_TEXT_MARGIN = 0.08
HEADER_TO_ARTIFACT = 0.12
ARTIFACT_TO_ARTIFACT = 0.18
ZONE_TOP_OFFSET = 0.08
HEADER_HEIGHT = 0.30
INSIGHT_LEFT_PADDING = 0.15
INSIGHT_RIGHT_PADDING = 0.12
INSIGHT_TOP_PADDING = 0.10
INSIGHT_BOTTOM_PADDING = 0.10
CARD_GAP = 0.15
CARD_INNER_PADDING = 0.15
TITLE_TO_SUBTITLE = 0.08
SUBTITLE_TO_BODY = 0.10
TABLE_MIN_ROW_HEIGHT = 0.32
CELL_PADDING = 0.09


def renderer_fallback_allowed(spec):
    """Whether Agent 6 is allowed to make design-time fallback decisions."""
    if not isinstance(spec, dict):
        return False
    policy = spec.get('fallback_policy') or {}
    return bool(
        spec.get('allow_renderer_fallback') or
        spec.get('allow_renderer_fit') or
        spec.get('allow_renderer_layout_fallback') or
        (isinstance(policy, dict) and (
            policy.get('allow_renderer_fallback') or
            policy.get('allow_renderer_fit') or
            policy.get('allow_renderer_layout_fallback')
        ))
    )


def fallback_policy_mode(spec, artifact_type=None, artifact_subtype=None):
    """Subtype-aware fallback mode supplied by Agent 5 block metadata."""
    if not isinstance(spec, dict):
        return None
    policy = spec.get('fallback_policy') or {}
    if not isinstance(policy, dict) or not renderer_fallback_allowed(spec):
        return None
    p_type = policy.get('artifact_type')
    p_subtype = policy.get('artifact_subtype')
    if artifact_type and p_type and p_type != artifact_type:
        return None
    if artifact_subtype and p_subtype and p_subtype != artifact_subtype:
        return None
    return policy.get('fallback_mode') or 'subtype_default'


def apply_block_render_metadata(artifact, block, default_type=None, default_subtype=None):
    """Copy Agent 5 block-level render metadata onto artifact dicts used by legacy renderers."""
    artifact = dict(artifact or {})
    if not isinstance(block, dict):
        return artifact
    artifact_type = block.get('artifact_type') or artifact.get('artifact_type') or default_type
    artifact_subtype = block.get('artifact_subtype') or artifact.get('artifact_subtype') or default_subtype
    if artifact_type:
        artifact['artifact_type'] = artifact_type
    if artifact_subtype:
        artifact['artifact_subtype'] = artifact_subtype
    if block.get('artifact_header_text') is not None:
        artifact['artifact_header_text'] = block.get('artifact_header_text')
    if block.get('block_role') is not None:
        artifact['block_role'] = block.get('block_role')
    policy = block.get('fallback_policy')
    if isinstance(policy, dict):
        artifact['fallback_policy'] = dict(policy)
        if policy.get('allow_renderer_fallback'):
            artifact['allow_renderer_fallback'] = True
        if policy.get('allow_renderer_fit'):
            artifact['allow_renderer_fit'] = True
        if policy.get('allow_renderer_layout_fallback'):
            artifact['allow_renderer_layout_fallback'] = True
    return artifact
CHART_HEADER_GAP = 0.11
OPTICAL_NUDGE = 0.02


# ─── TEMPLATE HELPERS ─────────────────────────────────────────────────────────

def clear_slides(prs):
    """Remove all content slides, preserving slide masters and layouts."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(nsmap.qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def find_blank_layout(prs):
    """Return the most minimal slide layout — prefers 'Blank', falls back by placeholder count."""
    for layout in prs.slide_layouts:
        if (layout.name or '').lower().strip() == 'blank':
            return layout
    for layout in prs.slide_layouts:
        if 'blank' in (layout.name or '').lower():
            return layout
    # Last resort: layout with fewest placeholders
    try:
        return min(prs.slide_layouts, key=lambda l: len(list(l.placeholders)))
    except Exception:
        return prs.slide_layouts[0]


def _is_non_content_layout(layout):
    """True for title/divider/closing layouts that should not back content slides."""
    try:
        lname = (layout.name or '').lower().strip()
    except Exception:
        lname = ''
    try:
        ltype = (layout._element.get('type', '') or '').lower().strip()
    except Exception:
        ltype = ''
    if ltype in ('title', 'sechead'):
        return True
    non_content_terms = (
        'title slide',
        'thank you',
        'closing',
        'end slide',
        'divider',
        'section header',
    )
    return any(term in lname for term in non_content_terms) or lname == 'section'


def find_content_fallback_layout(prs):
    """
    Pick a neutral layout for content slides when Agent 5 is in scratch mode.
    Never use title/divider/thank-you layouts as the base.
    """
    safe_layouts = [lyt for lyt in prs.slide_layouts if not _is_non_content_layout(lyt)]
    if not safe_layouts:
        return find_blank_layout(prs)

    preferred_names = ('blank', 'content', 'body text', 'topic', 'title and content')
    for pref in preferred_names:
        for lyt in safe_layouts:
            lname = (lyt.name or '').lower()
            if pref in lname:
                return lyt

    try:
        return min(safe_layouts, key=lambda l: len(list(l.placeholders)))
    except Exception:
        return safe_layouts[0]


def load_template_prs(template_b64):
    """
    Decode a base64 PPTX template, strip all content slides, and return the
    Presentation object. Slide masters and layouts are preserved so the brand
    master (background, fonts, logo, decorative shapes) carries through to
    every new slide.
    """
    template_bytes = base64.b64decode(template_b64)
    prs = Presentation(io.BytesIO(template_bytes))
    clear_slides(prs)
    return prs


# ─── HELPERS ──────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_color):
    """Convert #RRGGBB or #RGB to RGBColor."""
    if not hex_color or not isinstance(hex_color, str):
        return RGBColor(0x11, 0x11, 0x11)
    h = hex_color.lstrip('#')
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    if len(h) != 6:
        return RGBColor(0x11, 0x11, 0x11)
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return RGBColor(0x11, 0x11, 0x11)

def inches(val):
    """Safe inches conversion — handles None and strings."""
    try:
        return Inches(float(val or 0))
    except Exception:
        return Inches(0)

def pt(val):
    """Safe Pt conversion."""
    try:
        return Pt(float(val or 12))
    except Exception:
        return Pt(12)

def set_font(run_or_para, font_family, font_size, bold=False, italic=False, color_hex=None):
    """Apply font properties to a run or paragraph."""
    try:
        font = run_or_para.font
        if font_family:
            font.name = str(font_family)
        if font_size:
            font.size = pt(font_size)
        font.bold   = bool(bold)
        font.italic = bool(italic)
        if color_hex:
            font.color.rgb = hex_to_rgb(color_hex)
    except Exception:
        pass


def enable_text_fit(text_frame):
    """Keep text inside the shape whenever possible."""
    try:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
    except Exception:
        pass


def estimate_fit_font_size(text, width_in, height_in, max_size=12, min_size=7):
    """Cheap heuristic for table cells and dense text."""
    text = str(text or '')
    if not text.strip():
        return max_size

    width_pts = max(1, width_in * 72)
    height_pts = max(1, height_in * 72)
    usable_chars = max(6, int(width_pts / 5.4))
    lines = 0
    for chunk in text.split('\n'):
        lines += max(1, int(len(chunk) / usable_chars) + 1)

    if lines <= 1:
        return max_size

    size = min(max_size, int(height_pts / max(lines * 1.45, 1)))
    return max(min_size, size)


def estimate_wrapped_lines(text, width_in, font_size):
    """Cheap line-wrap estimate for short headings and labels."""
    text = str(text or '').strip()
    if not text:
        return 1
    width_pts = max(1, width_in * 72)
    char_w = max(1.0, font_size * 0.52)
    chars_per_line = max(4, int(width_pts / char_w))
    lines = 0
    for chunk in text.split('\n'):
        chunk = chunk.strip()
        lines += max(1, int(len(chunk) / chars_per_line) + 1)
    return max(1, lines)


def estimate_header_block_height(text, width_in, font_size):
    """Estimate rendered height for an artifact header, allowing wrapping."""
    lines = estimate_wrapped_lines(text, width_in, font_size)
    line_h = (font_size / 72.0) * 1.18
    return max(HEADER_HEIGHT, lines * line_h + 0.02)


def infer_slide_header_style(slide_spec):
    """Choose one consistent header language per content slide."""
    # Keep artifact headers visually consistent across all content slides.
    # Agent 5 controls placement; Agent 6 should not switch header language
    # slide-by-slide based on artifact sentiment.
    return 'underline'


def infer_artifact_header_style(artifact_type):
    """Choose header emphasis by artifact type."""
    if str(artifact_type or '').lower() == 'insight_text':
        return 'brand_fill'
    return 'underline'


def add_image_box(slide, image_b64, x, y, w, h):
    """Render an image from base64."""
    if not image_b64:
        return None
    try:
        blob = base64.b64decode(image_b64)
        return slide.shapes.add_picture(io.BytesIO(blob), inches(x), inches(y), inches(w), inches(h))
    except Exception:
        return None

def add_text_box(slide, x, y, w, h, text, font_family='Arial', font_size=12,
                 bold=False, color_hex='#111111', align='left', valign='top',
                 wrap=True, italic=False, hanging_in=0.0):
    """Add a text box to a slide with full formatting.

    hanging_in — if > 0, applies a hanging indent so that wrapped lines align
                 with the character after the bullet marker.
    """
    txBox = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    enable_text_fit(tf)

    # Vertical alignment
    valign_map = { 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM, 'top': MSO_ANCHOR.TOP }
    tf.vertical_anchor = valign_map.get(valign, MSO_ANCHOR.TOP)

    # Clear default paragraph and set text
    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text or '')

    # Paragraph alignment
    align_map = { 'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT }
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)

    set_font(run, font_family, font_size, bold, italic, color_hex)

    # Hanging indent: all continuation lines align with text after bullet
    if hanging_in > 0:
        _mar = int(Emu(inches(hanging_in)))
        pPr  = p._p.get_or_add_pPr()
        pPr.set('marL',   str(_mar))
        pPr.set('indent', str(-_mar))

    return txBox

def add_filled_rect(slide, x, y, w, h, fill_hex=None, border_hex=None,
                    border_pt=0, corner_radius=0):
    """Add a filled rectangle shape (used for backgrounds, cards, etc.)"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        inches(x), inches(y), inches(w), inches(h)
    )

    # Fill
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    else:
        shape.fill.background()

    # Border
    if border_hex and border_pt and border_pt > 0:
        shape.line.color.rgb = hex_to_rgb(border_hex)
        shape.line.width     = pt(border_pt)
    else:
        shape.line.fill.background()

    # Corner radius — set via XML
    if corner_radius and corner_radius > 0:
        try:
            sp  = shape._element
            prstGeom = sp.find('.//' + '{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prstGeom is not None:
                prstGeom.set('prst', 'roundRect')
                avLst = prstGeom.find('{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, '{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                # corner radius as percentage of shape size (max 50000 = 50%)
                radius_pct = min(int(corner_radius * 5000), 50000)
                gd = etree.SubElement(avLst, '{http://schemas.openxmlformats.org/drawingml/2006/main}gd')
                gd.set('name', 'adj')
                gd.set('fmla', 'val ' + str(radius_pct))
        except Exception:
            pass

    return shape

def rgb_tuple(hex_color):
    """Return (r, g, b) tuple from hex."""
    h = (hex_color or '#000000').lstrip('#')
    if len(h) == 3: h = h[0]*2 + h[1]*2 + h[2]*2
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return (0, 0, 0)


def is_dark_color(hex_color):
    """Return True if perceived luminance is < 0.5 (i.e. color is dark)."""
    r, g, b = rgb_tuple(hex_color)
    return (0.299 * r + 0.587 * g + 0.114 * b) / 255 < 0.5


def find_layout_by_name(prs, name):
    """Find a slide layout by exact name, then partial match."""
    if not name:
        return None
    name_lower = name.lower().strip()
    for layout in prs.slide_layouts:
        if (layout.name or '').lower().strip() == name_lower:
            return layout
    for layout in prs.slide_layouts:
        if name_lower in (layout.name or '').lower():
            return layout
    return None


def _compact_title_placeholder(slide, title_ph, text, font_size):
    """Shrink the title placeholder to fit its actual wrapped text height,
    then move all content placeholders up and expand them to reclaim the freed space.
    This eliminates the blank gap between a short title and the content area."""
    try:
        ph_w_in  = title_ph.width  / 914400.0   # EMU → inches
        n_lines  = estimate_wrapped_lines(text, max(0.5, ph_w_in), font_size)
        line_h   = (float(font_size) * 1.35) / 72.0   # inches per line incl. leading
        new_h_in = max(0.25, n_lines * line_h + 0.12)  # small top/bottom padding

        old_h_emu     = title_ph.height
        new_h_emu     = int(new_h_in * 914400)
        shrink_emu    = old_h_emu - new_h_emu         # positive = placeholder got shorter
        old_bottom    = title_ph.top + old_h_emu

        if shrink_emu <= int(0.05 * 914400):          # less than 0.05" — not worth moving
            return 0.0

        title_ph.height = new_h_emu

        # Move every non-title placeholder that sits at or below the old title bottom
        # upward by the full shrink amount and expand its height to keep the same bottom.
        threshold = old_bottom - int(0.3 * 914400)    # allow 0.3" tolerance
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                continue
            if ph.top >= threshold:
                ph.top    = max(0, ph.top - shrink_emu)
                ph.height = ph.height + shrink_emu

        return shrink_emu / 914400.0   # return shrink in inches so caller can shift free shapes
    except Exception as e:
        print('_compact_title_placeholder error:', e)
        return 0.0


def place_in_placeholder(slide, ph_idx, text, style_spec, bt,
                         preserve_template_style=False,
                         compact_title=True):
    """
    Write text into the placeholder at ph_idx on the slide.
    Falls back to a free-form text box if the placeholder is not found.
    For the title placeholder (idx 0), the placeholder is compacted to the actual
    text height so no blank gap appears between the title and the content area.
    """
    if not text:
        return
    style_spec = style_spec or {}
    bt         = bt or {}

    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == ph_idx:
                tf = ph.text_frame
                tf.clear()
                tf.word_wrap = True
                # Reset any vertical text direction inherited from the template layout.
                # Some layouts (e.g. "2 Column") have the title placeholder's bodyPr
                # set to vert="vert" which makes the title render sideways on the slide.
                try:
                    body_pr = tf._txBody.find(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr'
                    )
                    if body_pr is not None:
                        body_pr.attrib.pop('vert', None)
                except Exception:
                    pass
                p   = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(text)
                if not preserve_template_style:
                    font_family = style_spec.get('font_family') or bt.get('title_font_family', 'Arial')
                    font_size   = style_spec.get('font_size', 18 if ph_idx == 0 else 14)
                    bold        = style_spec.get('font_weight', '') in ('bold', 'semibold')
                    color_hex   = style_spec.get('color') or bt.get('title_color', '#111111')
                    align       = style_spec.get('align', 'left')
                    set_font(run, font_family, font_size, bold, False, color_hex)
                    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
                    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
                if ph_idx == 0 and compact_title and not preserve_template_style:
                    font_size = style_spec.get('font_size', 18 if ph_idx == 0 else 14)
                    return _compact_title_placeholder(slide, ph, text, font_size) or 0.0
                return 0.0
    except Exception as e:
        print(f'place_in_placeholder({ph_idx}) lookup error:', e)

    # Fallback: free-form text box using default positions
    if ph_idx == 0:
        add_text_box(slide, 0.4, 0.15, 9.0, 0.8, text,
                     style_spec.get('font_family', bt.get('title_font_family', 'Arial')),
                     style_spec.get('font_size', 20),
                     style_spec.get('font_weight', 'bold') in ('bold', 'semibold'),
                     style_spec.get('color', bt.get('title_color', '#1A3C8F')),
                     style_spec.get('align', 'left'), 'middle')
    else:
        add_text_box(slide, 0.4, 1.0, 9.0, 0.5, text,
                     style_spec.get('font_family', bt.get('body_font_family', 'Arial')),
                     style_spec.get('font_size', 14),
                     False,
                     style_spec.get('color', bt.get('body_color', '#333333')),
                     style_spec.get('align', 'left'), 'middle')
    return 0.0


def render_header_block(slide, header_block, bt, header_style='underline'):
    """
    Render an artifact header label.
    style='underline' — text + thin brand-colour rule below.
    style='brand_fill' — filled rectangle with contrasting text.
    """
    if not header_block:
        return
    x     = header_block.get('x', 0)
    y     = header_block.get('y', 0)
    w     = header_block.get('w', 4)
    h     = header_block.get('h', 0.3)
    text  = header_block.get('text', '')
    style = header_style or 'underline'

    if not text:
        return

    primary   = bt.get('primary_color', '#1A3C8F')
    font_fam  = bt.get('title_font_family', 'Arial')
    font_size = header_block.get('font_size', 11)
    text_h = estimate_header_block_height(text, w, font_size)
    rule_gap = 0.02
    rule_h = 0.005

    if style == 'brand_fill':
        fill_color = header_block.get('accent_color') or primary
        add_filled_rect(slide, x, y, w, h, fill_hex=fill_color)
        text_color = '#FFFFFF' if is_dark_color(fill_color) else '#111111'
        add_text_box(slide, x + 0.08, y, w - 0.16, h,
                     text, font_fam, font_size, True, text_color, 'left', 'middle')
        return y + h
    else:
        # Underline style
        add_text_box(slide, x, y, w, text_h,
                     text, font_fam, font_size, True, primary, 'left', 'top')
        rule_y = y + text_h + rule_gap
        # Use a connector (true hairline) instead of a filled rect to avoid PPTX minimum-height clamping
        from pptx.enum.shapes import MSO_CONNECTOR
        try:
            rule_line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                inches(x), inches(rule_y), inches(x + w), inches(rule_y)
            )
            rule_line.line.color.rgb = hex_to_rgb(primary)
            rule_line.line.width = pt(0.5)
        except Exception:
            pass
        return rule_y + rule_h


# ─── SLIDE RENDERERS ──────────────────────────────────────────────────────────

def render_insight_text(slide, artifact, bt, suppress_heading=False):
    """Render an insight_text artifact.

    suppress_heading — when True, the heading was already written to a layout header
                       placeholder; skip the inline heading rendering.
    """
    x = artifact.get('x', 0)
    y = artifact.get('y', 0)
    w = artifact.get('w', 4)
    h = artifact.get('h', 2)

    style   = artifact.get('style', {})
    hs      = artifact.get('heading_style', {})
    bs      = artifact.get('body_style', {})
    heading = artifact.get('heading', '') if not suppress_heading else ''
    points  = artifact.get('points', [])

    # Background fill + border — always render the container box; only the inline
    # heading text is suppressed when suppress_heading=True (because the header_block
    # was already rendered above the artifact by render_artifact).
    fill   = style.get('fill_color')
    border = style.get('border_color')
    bw     = style.get('border_width', 0)
    cr     = style.get('corner_radius', 0)
    if fill or (border and bw):
        add_filled_rect(slide, x, y, w, h, fill_hex=fill,
                        border_hex=border, border_pt=bw, corner_radius=cr)

    # ── Heading & body layout ────────────────────────────────────────────────
    TOP_PAD   = INSIGHT_TOP_PADDING
    content_y = y + TOP_PAD

    # Collect heading params; rendering is deferred until after b_size is known
    _h_params = None
    if heading:
        _h_params = {
            'font':  hs.get('font_family', bt.get('title_font_family', 'Arial')),
            'size':  hs.get('font_size', 12),
            'color': hs.get('color', bt.get('primary_color', '#1A3C8F')),
            'bold':  hs.get('font_weight', 'bold') in ('bold', 'semibold'),
        }
        content_y = y + 0.40   # heading row + gap

    # ── Native OOXML bullet helpers (defined early — used in sizing estimate too) ──
    BULLET_HANG = 0.22   # inches — bullet hangs left; text wraps at this indent

    def _strip_marker(s):
        s = str(s)
        return s[1:].lstrip() if s[:1] in ('✓', '✗', '→', '•') else s

    # Body points — compute sizes first so heading can scale to match
    if points:
        b_font       = bs.get('font_family', bt.get('body_font_family', 'Arial'))
        b_size       = bs.get('font_size', 10)
        b_color      = bs.get('color', '#111111')
        list_style   = bs.get('list_style', 'bullet')
        indent_in    = float(bs.get('indent_inches', INSIGHT_LEFT_PADDING))
        vert_dist    = bs.get('vertical_distribution', '')
        space_bef_pt = max(1, int(bs.get('space_before_pt', 4)))
        allow_fit_fallback = renderer_fallback_allowed(bs) or bs.get('font_size') is None
        rounded_inset = 0.04 if float(cr or 0) >= 4 else 0.0
        bullet_safe_inset = 0.08 + rounded_inset
        pad_top = float(bs.get('padding_top', max(INSIGHT_TOP_PADDING, 0.10 + rounded_inset)))
        pad_bottom = float(bs.get('padding_bottom', max(INSIGHT_BOTTOM_PADDING, 0.10 + rounded_inset)))
        pad_left = float(bs.get('padding_left', 0.12 + rounded_inset))
        pad_right = float(bs.get('padding_right', 0.10 + rounded_inset))
        pad_left = max(pad_left, 0.14 + rounded_inset)

        body_y = content_y + pad_top
        body_h = max(0.1, h - (body_y - y) - pad_bottom)
        para_left_in = max(BULLET_HANG + bullet_safe_inset, indent_in + bullet_safe_inset)
        text_w = max(0.5, w - pad_left - pad_right - para_left_in)

        # Estimate wrapped line count per point at a given font size
        def _wrapped_lines(text, font_pt, avail_w_in):
            # Average character width ≈ 50% of em (conservative)
            char_w = (font_pt / 72.0) * 0.50
            chars_per_line = max(1, int(avail_w_in / char_w))
            stripped = _strip_marker(str(text))
            words = stripped.split()
            if not words:
                return 1
            lines, line_len = 1, 0
            for word in words:
                wlen = len(word)
                if line_len == 0:
                    line_len = wlen
                elif line_len + 1 + wlen <= chars_per_line:
                    line_len += 1 + wlen
                else:
                    lines += 1
                    line_len = wlen
            return lines
        # Estimate total rendered height using actual wrap count.
        def _total_h(font_pt, gap_pt):
            lh = (font_pt / 72.0) * 1.30   # line height with leading
            gh = gap_pt / 72.0
            total = 0
            for idx, point in enumerate(points):
                total += _wrapped_lines(point, font_pt, text_w) * lh
                if idx > 0:
                    total += gh
            return total

        if allow_fit_fallback:
            total_h_est = _total_h(b_size, space_bef_pt)
            if total_h_est > body_h * 0.95:
                # Binary search for the largest font_pt that still fits
                lo, hi = 7, b_size
                while hi - lo > 0.5:
                    mid = (lo + hi) / 2.0
                    if _total_h(mid, space_bef_pt) <= body_h * 0.95:
                        lo = mid
                    else:
                        hi = mid
                b_size = max(7, int(lo))

            # Inter-bullet gap: renderer may adjust only on guarded fallback
            total_lines = sum(_wrapped_lines(p, b_size, text_w) for p in points)
            if total_lines > 12:
                b_size = max(7, b_size - 1)
                space_bef_pt = max(2, space_bef_pt - 1)

            # Keep heading/body proportionate only on guarded fallback
            if _h_params:
                spec_body          = bs.get('font_size', 10)
                _h_params['size']  = max(10, round(_h_params['size'] * (b_size / max(spec_body, 1))))

        if vert_dist == 'spread' and len(points) > 1:
            used_h = _total_h(b_size, space_bef_pt)
            slack_pt = max(0, int(round((body_h - used_h) * 72)))
            if slack_pt > 0:
                space_bef_pt = min(48, space_bef_pt + max(1, int(slack_pt / max(len(points) - 1, 1))))

    # Render heading now (font size finalised)
    if _h_params:
        add_text_box(slide, x + INSIGHT_LEFT_PADDING, y + TOP_PAD, w - INSIGHT_LEFT_PADDING, 0.34,
                     str(heading), _h_params['font'], _h_params['size'],
                     _h_params['bold'], _h_params['color'], 'left', 'middle')

    # ── Grouped mode ─────────────────────────────────────────────────────────
    insight_mode = artifact.get('insight_mode', 'standard')
    groups       = artifact.get('groups') or []

    if insight_mode == 'grouped' and groups:
        ghs      = artifact.get('group_header_style', {}) or {}
        gbs      = artifact.get('group_bullet_box_style', {}) or {}
        bsty     = artifact.get('bullet_style', {}) or {}
        g_gap    = float(artifact.get('group_gap_in')    or 0.08)
        hb_gap   = float(artifact.get('header_to_box_gap_in') or 0.05)
        g_layout = artifact.get('group_layout', 'rows')

        b_font   = bsty.get('font_family', bt.get('body_font_family', 'Arial'))
        b_size   = int(bsty.get('font_size', 10) or 10)
        b_color  = bsty.get('color', '#000000')
        b_char   = bsty.get('char', '•')
        b_space  = max(1, int(bsty.get('space_before_pt', 4) or 4))
        b_indent = float(bsty.get('indent_inches', 0.10) or 0.10)
        GRP_HANG = 0.15   # bullet hang indent inside group boxes

        h_font   = ghs.get('font_family', bt.get('title_font_family', 'Arial'))
        h_size   = int(ghs.get('font_size', 11) or 11)
        h_color  = ghs.get('text_color', '#FFFFFF')
        h_bold   = ghs.get('font_weight', 'bold') in ('bold', 'semibold')
        h_fill   = ghs.get('fill_color', bt.get('primary_color', '#0078AE'))
        h_cr     = float(ghs.get('corner_radius', 4) or 0)

        box_fill = gbs.get('fill_color')
        box_bdr  = gbs.get('border_color')
        box_bw   = float(gbs.get('border_width', 0.75) or 0)
        box_cr   = float(gbs.get('corner_radius', 4) or 0)
        box_pad  = gbs.get('padding', {}) or {}
        bp_t     = float(box_pad.get('top',    0.08) or 0)
        bp_r     = float(box_pad.get('right',  0.10) or 0)
        bp_b     = float(box_pad.get('bottom', 0.08) or 0)
        bp_l     = float(box_pad.get('left',   0.10) or 0)

        n             = len(groups)
        total_bullets = max(1, sum(len(g.get('bullets') or []) for g in groups))
        header_block_h = content_y - y   # space already used by header_block / heading

        def _render_group_bullets(bx, by, bw, bh, bullets):
            if not bullets:
                return
            bx2 = bx + bp_l + b_indent
            by2 = by + bp_t
            bw2 = max(0.2, bw - bp_l - bp_r - b_indent - GRP_HANG)
            bh2 = max(0.1, bh - bp_t - bp_b)
            txBox = slide.shapes.add_textbox(
                inches(bx2), inches(by2), inches(bw2), inches(bh2))
            tf = txBox.text_frame
            if renderer_fallback_allowed(artifact):
                enable_text_fit(tf)
            _hang_emu = int(Emu(inches(GRP_HANG)))
            for bi, bullet in enumerate(bullets):
                if bi == 0:
                    p_obj = tf.paragraphs[0]
                    p_obj.space_before = pt(0)
                else:
                    p_obj = tf.add_paragraph()
                    p_obj.space_before = pt(b_space)
                pPr = p_obj._p.get_or_add_pPr()
                for tag in [nsmap.qn('a:buNone'), nsmap.qn('a:buChar'),
                            nsmap.qn('a:buAutoNum'), nsmap.qn('a:buFont')]:
                    for el in pPr.findall(tag):
                        pPr.remove(el)
                pPr.set('marL',   str(_hang_emu))
                pPr.set('indent', str(-_hang_emu))
                buFont = etree.SubElement(pPr, nsmap.qn('a:buFont'))
                buFont.set('typeface', b_font)
                buChar = etree.SubElement(pPr, nsmap.qn('a:buChar'))
                buChar.set('char', b_char)
                run = p_obj.add_run()
                run.text = _strip_marker(str(bullet))
                set_font(run, b_font, b_size, False, False, b_color)

        if g_layout == 'rows':
            h_w      = float(ghs.get('w', 1.2) or 1.2)
            box_x    = x + h_w + hb_gap
            box_w    = max(0.2, w - h_w - hb_gap - INSIGHT_RIGHT_PADDING)
            total_rh = max(0.2, h - header_block_h - INSIGHT_BOTTOM_PADDING
                           - (n - 1) * g_gap)
            cur_y    = content_y
            for g in groups:
                nbullets = max(1, len(g.get('bullets') or []))
                row_h    = max(0.25, total_rh * (nbullets / total_bullets))
                # Header label (left column)
                add_filled_rect(slide, x, cur_y, h_w, row_h,
                                fill_hex=h_fill, corner_radius=h_cr)
                add_text_box(slide, x + 0.06, cur_y, h_w - 0.12, row_h,
                             str(g.get('header', '')),
                             h_font, h_size, h_bold, h_color, 'center', 'middle')
                # Bullet box (right column)
                if box_fill or (box_bdr and box_bw):
                    add_filled_rect(slide, box_x, cur_y, box_w, row_h,
                                    fill_hex=box_fill, border_hex=box_bdr,
                                    border_pt=box_bw, corner_radius=box_cr)
                _render_group_bullets(box_x, cur_y, box_w, row_h,
                                      g.get('bullets') or [])
                cur_y += row_h + g_gap

        else:  # columns
            col_w   = max(0.2, (w - INSIGHT_RIGHT_PADDING - (n - 1) * g_gap) / max(n, 1))
            h_h     = float(ghs.get('h', 0.28) or 0.28)
            box_h   = max(0.1, h - header_block_h - h_h - hb_gap - INSIGHT_BOTTOM_PADDING)
            cur_x   = x
            for g in groups:
                # Header label (top of column)
                add_filled_rect(slide, cur_x, content_y, col_w, h_h,
                                fill_hex=h_fill, corner_radius=h_cr)
                add_text_box(slide, cur_x + 0.05, content_y, col_w - 0.10, h_h,
                             str(g.get('header', '')),
                             h_font, h_size, h_bold, h_color, 'center', 'middle')
                # Bullet box (below header)
                bullet_y = content_y + h_h + hb_gap
                if box_fill or (box_bdr and box_bw):
                    add_filled_rect(slide, cur_x, bullet_y, col_w, box_h,
                                    fill_hex=box_fill, border_hex=box_bdr,
                                    border_pt=box_bw, corner_radius=box_cr)
                _render_group_bullets(cur_x, bullet_y, col_w, box_h,
                                      g.get('bullets') or [])
                cur_x += col_w + g_gap
        return   # grouped mode fully rendered

    def _bullet_char(point, idx):
        s = str(point)
        if list_style == 'tick_cross':
            if s[:1] in ('✓', '✗', '→'):
                return s[:1]
            return '✗' if any(neg in s.lower() for neg in
                ['not ', 'no ', 'declin', 'risk', 'fail', 'loss', 'below', 'miss']) else '✓'
        if list_style == 'numbered':
            return str(idx + 1) + '.'
        return '•'

    def _set_bullet(pPr, bchar, hang_emu):
        """Apply native OOXML hanging-indent bullet to an <a:pPr> element."""
        for tag in [nsmap.qn('a:buNone'), nsmap.qn('a:buChar'),
                    nsmap.qn('a:buAutoNum'), nsmap.qn('a:buFont')]:
            for el in pPr.findall(tag):
                pPr.remove(el)
        pPr.set('marL',   str(hang_emu))
        pPr.set('indent', str(-hang_emu))
        buFont = etree.SubElement(pPr, nsmap.qn('a:buFont'))
        buFont.set('typeface', b_font)
        buChar = etree.SubElement(pPr, nsmap.qn('a:buChar'))
        buChar.set('char', bchar)

    if points:
        _hang_emu = int(Emu(inches(BULLET_HANG)))
        _mar_left_emu = int(Emu(inches(para_left_in)))

        # ALWAYS use a single text box for all points.
        # For "spread" distribution: use larger space_before on each paragraph so
        # points fill the available height — this keeps the text in one container
        # so the renderer can auto-fit the font if needed.
        txBox = slide.shapes.add_textbox(
            inches(x + pad_left), inches(body_y),
            inches(w - pad_left - pad_right), inches(body_h)
        )
        tf = txBox.text_frame
        if allow_fit_fallback:
            enable_text_fit(tf)
        try:
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
        except Exception:
            pass

        for i, point in enumerate(points):
            if i == 0:
                p_obj = tf.paragraphs[0]
                # First paragraph: no space before (avoid top gap)
                p_obj.space_before = pt(0)
            else:
                p_obj = tf.add_paragraph()
                p_obj.space_before = pt(space_bef_pt)
            pPr = p_obj._p.get_or_add_pPr()
            _set_bullet(pPr, _bullet_char(point, i), _hang_emu)
            pPr.set('marL', str(_mar_left_emu))
            run = p_obj.add_run()
            run.text = _strip_marker(point)
            set_font(run, b_font, b_size, False, False, b_color)


def _apply_dual_axis(chart, secondary_series_names):
    """Post-process a clustered column chart XML to move named series onto a secondary
    line-chart plot with its own right-hand Y axis.

    python-pptx does not natively support combo charts, so we manipulate the XML directly.
    After this call the chart has two plot elements: barChart (primary) + lineChart (secondary).
    """
    plot_area = chart._element.find(nsmap.qn('c:plotArea'))
    if plot_area is None:
        return

    bar_chart = plot_area.find(nsmap.qn('c:barChart'))
    if bar_chart is None:
        return

    # Grab existing axId values (catAx id, primary valAx id)
    ax_id_els = bar_chart.findall(nsmap.qn('c:axId'))
    if len(ax_id_els) < 2:
        return
    cat_ax_id     = ax_id_els[0].get('val')
    prim_val_ax_id = ax_id_els[1].get('val')
    sec_val_ax_id  = str(int(prim_val_ax_id) + 1)   # unique secondary ID

    # Pull matching series elements out of barChart
    sec_ser_els = []
    for ser_el in list(bar_chart.findall(nsmap.qn('c:ser'))):
        tx = ser_el.find('.//' + nsmap.qn('c:v'))
        if tx is not None and tx.text in secondary_series_names:
            bar_chart.remove(ser_el)
            sec_ser_els.append(ser_el)

    if not sec_ser_els:
        return   # nothing matched — leave chart unchanged

    # Re-index remaining bar series
    for idx, ser_el in enumerate(bar_chart.findall(nsmap.qn('c:ser'))):
        for tag in (nsmap.qn('c:idx'), nsmap.qn('c:order')):
            el = ser_el.find(tag)
            if el is not None:
                el.set('val', str(idx))

    # Build lineChart element for secondary series
    line_chart = etree.SubElement(plot_area, nsmap.qn('c:lineChart'))
    etree.SubElement(line_chart, nsmap.qn('c:grouping')).set('val', 'standard')
    etree.SubElement(line_chart, nsmap.qn('c:varyColors')).set('val', '0')

    prim_bar_count = len(bar_chart.findall(nsmap.qn('c:ser')))
    for idx, ser_el in enumerate(sec_ser_els):
        for tag in (nsmap.qn('c:idx'), nsmap.qn('c:order')):
            el = ser_el.find(tag)
            if el is not None:
                el.set('val', str(prim_bar_count + idx))
        line_chart.append(ser_el)

    etree.SubElement(line_chart, nsmap.qn('c:marker')).set('val', '1')
    etree.SubElement(line_chart, nsmap.qn('c:smooth')).set('val', '0')
    etree.SubElement(line_chart, nsmap.qn('c:axId')).set('val', cat_ax_id)
    etree.SubElement(line_chart, nsmap.qn('c:axId')).set('val', sec_val_ax_id)

    # Add secondary valAx after existing valAx
    sec_val_ax = etree.SubElement(plot_area, nsmap.qn('c:valAx'))
    etree.SubElement(sec_val_ax, nsmap.qn('c:axId')).set('val', sec_val_ax_id)
    scaling = etree.SubElement(sec_val_ax, nsmap.qn('c:scaling'))
    etree.SubElement(scaling, nsmap.qn('c:orientation')).set('val', 'minMax')
    etree.SubElement(sec_val_ax, nsmap.qn('c:delete')).set('val', '0')
    etree.SubElement(sec_val_ax, nsmap.qn('c:axPos')).set('val', 'r')
    etree.SubElement(sec_val_ax, nsmap.qn('c:crossAx')).set('val', cat_ax_id)
    etree.SubElement(sec_val_ax, nsmap.qn('c:crosses')).set('val', 'max')
    etree.SubElement(sec_val_ax, nsmap.qn('c:crossBetween')).set('val', 'between')


def _estimate_legend_text_width(label, font_size_pt):
    text = str(label or '')
    return max(0.40, min(2.20, len(text) * max(font_size_pt, 8) * 0.0105))


def _chart_legend_entries(chart_type_str, categories, series_data, series_styles, chart_palette, allow_fallback):
    entries = []
    # pie, donut, and group_pie: legend entries represent SLICES (categories), colored per series_style[i]
    if chart_type_str in ('pie', 'donut', 'group_pie'):
        for i, category in enumerate(categories or []):
            style = series_styles[i] if i < len(series_styles) else {}
            color = style.get('fill_color')
            if not color and allow_fallback:
                color = chart_palette[i % len(chart_palette)]
            entries.append({
                'label': str(category or ''),
                'color': color or '#666666',
            })
        return entries

    for i, series in enumerate(series_data or []):
        style = series_styles[i] if i < len(series_styles) else {}
        color = style.get('fill_color')
        if not color and allow_fallback:
            color = chart_palette[i % len(chart_palette)]
        entries.append({
            'label': str(series.get('name') or ('Series ' + str(i + 1))),
            'color': color or '#666666',
        })
    return entries


def _compute_chart_legend_layout(x, y, w, h, legend_position, legend_entries, font_size_pt):
    if not legend_entries or legend_position not in ('top', 'right'):
        return ((x, y, w, h), None)

    swatch = 0.14
    text_gap = 0.06
    item_gap_x = 0.18
    row_gap = 0.06
    line_h = max(0.20, font_size_pt * 0.022)
    pad_x = 0.04
    pad_y = 0.03

    if legend_position == 'top':
        rows = []
        current = []
        used_w = 0.0
        max_row_w = max(0.5, w - 0.04)
        for entry in legend_entries:
            item_w = swatch + text_gap + _estimate_legend_text_width(entry['label'], font_size_pt)
            proposed = item_w if not current else used_w + item_gap_x + item_w
            if current and proposed > max_row_w:
                rows.append(current)
                current = [dict(entry, item_w=item_w)]
                used_w = item_w
            else:
                current.append(dict(entry, item_w=item_w))
                used_w = proposed
        if current:
            rows.append(current)

        legend_h = pad_y * 2 + len(rows) * line_h + max(0, len(rows) - 1) * row_gap
        legend_h = min(max(legend_h, 0.28), h * 0.28)
        chart_y = y + legend_h + 0.05
        chart_h = max(1.0, h - legend_h - 0.05)
        legend_box = {
            'position': 'top',
            'x': x,
            'y': y,
            'w': w,
            'h': legend_h,
            'rows': rows,
            'line_h': line_h,
            'pad_x': pad_x,
            'pad_y': pad_y,
            'swatch': swatch,
            'text_gap': text_gap,
            'item_gap_x': item_gap_x,
            'row_gap': row_gap,
        }
        return ((x, chart_y, w, chart_h), legend_box)

    max_item_w = 0.0
    items = []
    for entry in legend_entries:
        item_w = swatch + text_gap + _estimate_legend_text_width(entry['label'], font_size_pt)
        max_item_w = max(max_item_w, item_w)
        items.append(dict(entry, item_w=item_w))
    legend_w = min(max(max_item_w + pad_x * 2, 1.05), w * 0.38)
    chart_w = max(1.0, w - legend_w - 0.08)
    legend_box = {
        'position': 'right',
        'x': x + chart_w + 0.08,
        'y': y,
        'w': legend_w,
        'h': h,
        'items': items,
        'line_h': line_h,
        'pad_x': pad_x,
        'pad_y': pad_y,
        'swatch': swatch,
        'text_gap': text_gap,
        'item_gap_x': item_gap_x,
        'row_gap': row_gap,
    }
    return ((x, y, chart_w, h), legend_box)


def _render_custom_chart_legend(slide, legend_box, font_family, font_size_pt, color_hex):
    if not legend_box:
        return

    if legend_box.get('position') == 'top':
        cur_y = legend_box['y'] + legend_box['pad_y']
        for row in legend_box.get('rows', []):
            row_w = sum(item['item_w'] for item in row) + max(0, len(row) - 1) * legend_box['item_gap_x']
            cur_x = legend_box['x'] + max(legend_box['pad_x'], (legend_box['w'] - row_w) / 2.0)
            for item in row:
                add_filled_rect(
                    slide,
                    cur_x,
                    cur_y + max(0.0, (legend_box['line_h'] - legend_box['swatch']) / 2.0),
                    legend_box['swatch'],
                    legend_box['swatch'],
                    fill_hex=item['color']
                )
                add_text_box(
                    slide,
                    cur_x + legend_box['swatch'] + legend_box['text_gap'],
                    cur_y,
                    max(0.35, item['item_w'] - legend_box['swatch'] - legend_box['text_gap']),
                    legend_box['line_h'],
                    item['label'],
                    font_family=font_family,
                    font_size=font_size_pt,
                    bold=False,
                    color_hex=color_hex,
                    align='left',
                    valign='middle',
                    wrap=False
                )
                cur_x += item['item_w'] + legend_box['item_gap_x']
            cur_y += legend_box['line_h'] + legend_box['row_gap']
        return

    cur_y = legend_box['y'] + legend_box['pad_y']
    for item in legend_box.get('items', []):
        add_filled_rect(
            slide,
            legend_box['x'] + legend_box['pad_x'],
            cur_y + max(0.0, (legend_box['line_h'] - legend_box['swatch']) / 2.0),
            legend_box['swatch'],
            legend_box['swatch'],
            fill_hex=item['color']
        )
        add_text_box(
            slide,
            legend_box['x'] + legend_box['pad_x'] + legend_box['swatch'] + legend_box['text_gap'],
            cur_y,
            max(0.35, legend_box['w'] - (legend_box['pad_x'] * 2) - legend_box['swatch'] - legend_box['text_gap']),
            legend_box['line_h'],
            item['label'],
            font_family=font_family,
            font_size=font_size_pt,
            bold=False,
            color_hex=color_hex,
            align='left',
            valign='middle',
            wrap=True
        )
        cur_y += legend_box['line_h'] + legend_box['row_gap']


def render_group_pie(slide, x, y, w, h, categories, series_data, series_styles,
                     chart_palette, show_labels, label_size, accent_color, body_font,
                     body_color='#444444'):
    """Render N individual pie charts side-by-side in a horizontal group.

    x, y, w, h       — bounding box AFTER legend space has been reserved above
    categories        — shared slice labels (same for every pie)
    series_data       — one entry per entity/pie: {'name': str, 'series_total': str, 'values': [float,...]}
    series_styles     — one entry per SLICE (category): {'fill_color': hex, ...}
    accent_color      — brand accent used for entity name labels below each pie
    body_font         — font family for entity labels
    body_color        — secondary text color used for series_total sub-labels
    """
    n_pies = len(series_data)
    if n_pies == 0 or not categories:
        return

    # Determine whether any series carries a series_total sub-label
    has_any_total   = any(str(s.get('series_total') or '').strip() for s in series_data)
    ENTITY_LABEL_H  = 0.44 if has_any_total else 0.28   # taller when sub-label present
    ENTITY_LABEL_GAP = 0.05   # gap between pie bottom and label top

    pie_col_w = w / n_pies                       # width allocated to each pie column
    pie_area_h = h - ENTITY_LABEL_H - ENTITY_LABEL_GAP  # height available for the pie circle
    pie_size  = max(0.5, min(pie_col_w * 0.90, pie_area_h))  # square, 90% of column width

    for i, ser in enumerate(series_data):
        entity_name = str(ser.get('name') or '')
        values = [float(v) if v is not None else 0.0 for v in (ser.get('values') or [])]
        while len(values) < len(categories):
            values.append(0.0)
        values = values[:len(categories)]

        # Centre the pie within its column
        col_left    = x + i * pie_col_w
        pie_x       = col_left + (pie_col_w - pie_size) / 2
        pie_y       = y + max(0, (pie_area_h - pie_size) / 2)

        # Build single-series ChartData for this pie
        cd = ChartData()
        cd.categories = [str(c) for c in categories]
        cd.add_series(entity_name, values)

        try:
            pie_shape = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE,
                inches(pie_x), inches(pie_y), inches(pie_size), inches(pie_size),
                cd
            )
            pie_chart = pie_shape.chart
            pie_chart.has_title  = False
            pie_chart.has_legend = False

            # Color each slice from series_styles (index = slice/category index)
            try:
                for ser_obj in pie_chart.series:
                    for pi, point_obj in enumerate(ser_obj.points):
                        c_hex = (series_styles[pi].get('fill_color')
                                 if pi < len(series_styles) else None)
                        if not c_hex:
                            c_hex = chart_palette[pi % max(len(chart_palette), 1)]
                        if c_hex:
                            point_obj.format.fill.solid()
                            point_obj.format.fill.fore_color.rgb = hex_to_rgb(c_hex)
            except Exception:
                pass

            # Percentage data labels — built entirely via raw XML so element ordering
            # is correct per OOXML spec:
            #   <c:dLbl> delete overrides FIRST, then numFmt/txPr/position/show flags.
            # Using the python-pptx API alone reorders elements and breaks both the
            # zero-label suppression and the dLblPos override.
            if show_labels:
                try:
                    for ser_obj in pie_chart.series:
                        ser_el = ser_obj._element

                        # Remove any <c:dLbls> python-pptx may have already written
                        for existing in ser_el.findall(nsmap.qn('c:dLbls')):
                            ser_el.remove(existing)

                        dLbls_el = etree.Element(nsmap.qn('c:dLbls'))

                        # ── 1. Individual label overrides (MUST be first in dLbls) ──
                        # Delete labels whose slice value is 0 so "0%" never appears
                        for pi, val in enumerate(values):
                            if val == 0.0:
                                dLbl = etree.SubElement(dLbls_el, nsmap.qn('c:dLbl'))
                                etree.SubElement(dLbl, nsmap.qn('c:idx')).set('val', str(pi))
                                etree.SubElement(dLbl, nsmap.qn('c:delete')).set('val', '1')

                        # ── 2. Number format ──
                        numFmt = etree.SubElement(dLbls_el, nsmap.qn('c:numFmt'))
                        numFmt.set('formatCode', '0%')
                        numFmt.set('sourceLinked', '0')

                        # ── 3. Font color + size via txPr ──
                        txPr = etree.SubElement(dLbls_el, nsmap.qn('c:txPr'))
                        etree.SubElement(txPr, nsmap.qn('a:bodyPr'))
                        etree.SubElement(txPr, nsmap.qn('a:lstStyle'))
                        p_el   = etree.SubElement(txPr, nsmap.qn('a:p'))
                        pPr_el = etree.SubElement(p_el, nsmap.qn('a:pPr'))
                        defRPr = etree.SubElement(pPr_el, nsmap.qn('a:defRPr'))
                        defRPr.set('sz', str(int(max(7, label_size) * 100)))  # hundredths of a pt
                        solidFill = etree.SubElement(defRPr, nsmap.qn('a:solidFill'))
                        srgbClr   = etree.SubElement(solidFill, nsmap.qn('a:srgbClr'))
                        srgbClr.set('val', accent_color.lstrip('#'))

                        # ── 4. Label position — always outside end ──
                        etree.SubElement(dLbls_el, nsmap.qn('c:dLblPos')).set('val', 'outEnd')

                        # ── 5. Leader lines for small/overlapping slices (e.g. Kirana) ──
                        etree.SubElement(dLbls_el, nsmap.qn('c:showLdrLines')).set('val', '1')

                        # ── 6. Show flags ──
                        etree.SubElement(dLbls_el, nsmap.qn('c:showLegendKey')).set('val', '0')
                        etree.SubElement(dLbls_el, nsmap.qn('c:showVal')).set('val', '0')
                        etree.SubElement(dLbls_el, nsmap.qn('c:showCatName')).set('val', '0')
                        etree.SubElement(dLbls_el, nsmap.qn('c:showSerName')).set('val', '0')
                        etree.SubElement(dLbls_el, nsmap.qn('c:showPercent')).set('val', '1')
                        etree.SubElement(dLbls_el, nsmap.qn('c:showBubbleSize')).set('val', '0')

                        # Insert dLbls before <c:cat> to keep OOXML element order valid
                        cat_el = ser_el.find(nsmap.qn('c:cat'))
                        if cat_el is not None:
                            ser_el.insert(list(ser_el).index(cat_el), dLbls_el)
                        else:
                            ser_el.append(dLbls_el)
                except Exception as e:
                    print(f'render_group_pie: data label XML error:', e)
        except Exception as e:
            print(f'render_group_pie: error rendering pie {i} ({entity_name}):', e)

        # Entity label below the pie — centre-aligned, brand accent colour
        label_y      = y + pie_area_h + ENTITY_LABEL_GAP
        series_total = str(ser.get('series_total') or '').strip()

        if series_total:
            # Two-line block: entity name (bold, accent) + series_total (regular, body)
            NAME_H  = 0.24
            TOTAL_H = 0.18
            add_text_box(
                slide,
                col_left, label_y, pie_col_w, NAME_H,
                entity_name,
                body_font, 10, True, accent_color,
                'center', 'middle'
            )
            add_text_box(
                slide,
                col_left, label_y + NAME_H, pie_col_w, TOTAL_H,
                series_total,
                body_font, 8, False, body_color,
                'center', 'top'
            )
        else:
            add_text_box(
                slide,
                col_left, label_y, pie_col_w, ENTITY_LABEL_H,
                entity_name,
                body_font, 10, True, accent_color,
                'center', 'middle'
            )


def render_chart(slide, artifact, bt, suppress_heading=False, slide_w=13.33, slide_h=7.5):
    """Render a chart artifact using python-pptx native charts.

    suppress_heading — when True the artifact heading was already written to a layout
                       placeholder; the internal chart title is suppressed.
    """
    x  = artifact.get('x', 0)
    y  = artifact.get('y', 0)
    w  = artifact.get('w', 5)
    h  = artifact.get('h', 3)

    chart_type_str  = artifact.get('chart_type', 'bar')
    categories      = artifact.get('categories', [])
    series_data     = artifact.get('series', [])
    chart_title     = artifact.get('chart_title', '')
    show_labels     = artifact.get('show_data_labels', True)
    show_legend     = artifact.get('show_legend', False)
    series_styles   = artifact.get('series_style', [])
    cs              = artifact.get('chart_style', {})
    dual_axis       = artifact.get('dual_axis', False)
    secondary_names = set(artifact.get('secondary_series', []))
    chart_palette   = bt.get('chart_palette', ['#0F2FB5', '#FF8E00', '#2D962D', '#D60202'])
    header_block    = artifact.get('header_block', {}) or {}
    header_font_size = int(header_block.get('font_size') or cs.get('title_font_size', 11) or 11)

    # Read Agent-5-pre-computed fields (if present) for legend, label size, rotation.
    art_computed = artifact.get('_computed', {}) or {}
    _computed_legend_pos   = art_computed.get('legend_position', None)
    _computed_label_size   = art_computed.get('data_label_size', None)
    _computed_cat_rotation = art_computed.get('category_label_rotation', None)
    chart_subtype = artifact.get('artifact_subtype') or chart_type_str
    chart_fallback_mode = fallback_policy_mode(artifact, 'chart', chart_subtype)
    allow_chart_fallback = bool(chart_fallback_mode)

    # Adaptive data label size — always a proportion of header_font_size.
    # A density factor (0.0–1.0) scales down the ratio when categories are
    # tightly packed relative to the chart's available space.
    # Use Agent 5 pre-computed value when available, else fall back to heuristic.
    if _computed_label_size is not None:
        max_chart_label_size = int(_computed_label_size)
    elif allow_chart_fallback:
        _n_cat   = max(1, len(categories))
        _chart_w = float(w or 5)
        _chart_h = float(h or 3)
        if chart_type_str == 'horizontal_bar':
            # Space per bar in inches; 0.55" is a comfortable single-bar baseline
            _density = min(1.0, (_chart_h / _n_cat) / 0.55)
        elif chart_type_str in ('bar', 'clustered_bar', 'line', 'waterfall'):
            # Space per column in inches; 0.65" is a comfortable single-column baseline
            _density = min(1.0, (_chart_w / _n_cat) / 0.65)
        else:
            _density = 1.0
        # At full density: labels at 75% of header. Scales down proportionally with density.
        # Floor: never below 55% of header (keeps labels readable even on dense charts).
        max_chart_label_size = max(
            round(header_font_size * 0.55),
            round(header_font_size * 0.75 * _density)
        )
    else:
        max_chart_label_size = int(cs.get('label_font_size') or header_font_size or 9)

    # Map chart type string to XL_CHART_TYPE
    chart_type_map = {
        'bar':           XL_CHART_TYPE.COLUMN_CLUSTERED,
        'clustered_bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'horizontal_bar':XL_CHART_TYPE.BAR_CLUSTERED,
        'line':          XL_CHART_TYPE.LINE,
        'pie':           XL_CHART_TYPE.PIE,
        'waterfall':     XL_CHART_TYPE.COLUMN_CLUSTERED,
    }
    xl_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    if not categories or not series_data:
        add_filled_rect(slide, x, y, w, h, fill_hex='#F0F4FF', border_hex='#C0C8E0', border_pt=1)
        add_text_box(slide, x + 0.1, y + h / 2 - 0.2, w - 0.2, 0.4,
                     'Chart: ' + chart_type_str + ' (no data)', 'Arial', 10,
                     False, '#888888', 'center', 'middle')
        return

    _legend_pos_map = {'right': 4, 'top': 1, 'bottom': 3, 'left': 2, 'none': None}
    if _computed_legend_pos is not None and _computed_legend_pos in _legend_pos_map:
        resolved_legend_position = _computed_legend_pos
        legend_pos = _legend_pos_map[_computed_legend_pos]
    elif allow_chart_fallback:
        chart_w_ratio = (float(w) / float(slide_w)) if slide_w > 0 else 0.0
        chart_h_ratio = (float(h) / float(slide_h)) if slide_h > 0 else 0.0
        if chart_h_ratio > 0.60:
            resolved_legend_position = 'top'
            legend_pos = 1
        elif chart_w_ratio > 0.60:
            resolved_legend_position = 'right'
            legend_pos = 4
        elif chart_type_str == 'pie':
            resolved_legend_position = 'right'
            legend_pos = 4
        else:
            resolved_legend_position = 'top'
            legend_pos = 1
    else:
        resolved_legend_position = str(cs.get('legend_position', 'none') or 'none').lower()
        legend_pos = _legend_pos_map.get(resolved_legend_position)

    _effective_show_legend = show_legend and (legend_pos is not None)
    legend_font_size = int(cs.get('legend_font_size', header_font_size) or header_font_size)
    legend_font_size = min(legend_font_size, max(8, header_font_size - 1), 9)
    legend_entries = _chart_legend_entries(
        chart_type_str, categories, series_data, series_styles, chart_palette, allow_chart_fallback
    ) if _effective_show_legend else []
    (chart_x, chart_y, chart_w, chart_h), custom_legend_box = _compute_chart_legend_layout(
        float(x), float(y), float(w), float(h), resolved_legend_position, legend_entries, legend_font_size
    )

    # ── group_pie: delegate to dedicated renderer, then render legend and return ──
    if chart_type_str == 'group_pie':
        if _effective_show_legend and custom_legend_box:
            _render_custom_chart_legend(
                slide, custom_legend_box,
                cs.get('legend_font_family', bt.get('body_font_family', 'Arial')),
                legend_font_size,
                cs.get('legend_color', bt.get('body_color', '#000000'))
            )
        accent_color = next(
            (s.get('data_label_color') for s in series_styles if s.get('data_label_color')),
            None
        ) or bt.get('body_color') or bt.get('primary_color') or '#1A1A1A'
        body_font    = cs.get('label_font_family') or bt.get('body_font_family', 'Arial')
        body_color   = bt.get('body_color') or bt.get('secondary_color') or '#444444'
        render_group_pie(
            slide,
            chart_x, chart_y, chart_w, chart_h,
            categories, series_data, series_styles,
            chart_palette, show_labels, max_chart_label_size,
            accent_color, body_font, body_color
        )
        return

    # Build ChartData
    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    for si, ser in enumerate(series_data):
        ser_name   = ser.get('name', 'Series ' + str(si + 1))
        ser_values = [float(v) if v is not None else 0 for v in (ser.get('values') or [])]
        while len(ser_values) < len(categories): ser_values.append(0)
        cd.add_series(ser_name, ser_values[:len(categories)])

    # Add chart
    chart = slide.shapes.add_chart(
        xl_type, inches(chart_x), inches(chart_y), inches(chart_w), inches(chart_h), cd
    ).chart

    # Apply dual-axis layout (barChart primary + lineChart secondary with right Y axis)
    if dual_axis and secondary_names and chart_type_str not in ('pie', 'line'):
        try:
            _apply_dual_axis(chart, secondary_names)
        except Exception as e:
            print('dual axis error:', e)

    # Title — suppress when the zone header placeholder already carries the label
    show_title = bool(chart_title) and not suppress_heading
    chart.has_title = show_title
    if show_title:
        chart.chart_title.text_frame.text = str(chart_title)
        try:
            tf_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
            set_font(tf_run,
                     cs.get('title_font_family', 'Arial'),
                     cs.get('title_font_size', 12),
                     True, False,
                     cs.get('title_color', '#000000'))
        except Exception:
            pass

    # Legend placement — prefer Agent 5 pre-computed value; fall back to heuristic.
    # XL_LEGEND_POSITION: BOTTOM=3  RIGHT=4  TOP=1  LEFT=2
    _legend_pos_map = {'right': 4, 'top': 1, 'bottom': 3, 'left': 2, 'none': None}
    if _computed_legend_pos is not None and _computed_legend_pos in _legend_pos_map:
        legend_pos = _legend_pos_map[_computed_legend_pos]
    elif allow_chart_fallback:
        # Heuristic:
        #   a) chart height > 60% of actual slide height → TOP
        #   b) else if chart width  > 60% of actual slide width  → RIGHT
        #   c) else pie charts → RIGHT; other charts → TOP
        chart_w_ratio = (float(w) / float(slide_w)) if slide_w > 0 else 0.0
        chart_h_ratio = (float(h) / float(slide_h)) if slide_h > 0 else 0.0
        if chart_h_ratio > 0.60:
            legend_pos = 1   # TOP
        elif chart_w_ratio > 0.60:
            legend_pos = 4   # RIGHT
        elif chart_type_str == 'pie':
            legend_pos = 4   # RIGHT
        else:
            legend_pos = 1   # TOP
    else:
        legend_pos = _legend_pos_map.get(cs.get('legend_position', 'none'))

    # If _computed says 'none', treat as no legend regardless of show_legend flag
    _effective_show_legend = show_legend and (legend_pos is not None)
    chart.has_legend = False

    # Gridlines are always suppressed.
    try:
        if hasattr(chart, 'value_axis') and chart.value_axis is not None:
            chart.value_axis.has_major_gridlines = False
            chart.value_axis.has_minor_gridlines = False
    except Exception:
        pass

    if _effective_show_legend and custom_legend_box:
        _render_custom_chart_legend(
            slide,
            custom_legend_box,
            cs.get('legend_font_family', bt.get('body_font_family', 'Arial')),
            legend_font_size,
            cs.get('legend_color', bt.get('body_color', '#000000'))
        )

    # Series colors + data labels
    try:
        for si, ser_obj in enumerate(chart.series):
            color_hex = None
            if si < len(series_styles):
                color_hex = series_styles[si].get('fill_color')
            if not color_hex and allow_chart_fallback:
                color_hex = chart_palette[si % len(chart_palette)]
            try:
                if color_hex:
                    ser_obj.format.fill.solid()
                    ser_obj.format.fill.fore_color.rgb = hex_to_rgb(color_hex)
            except Exception:
                pass
            if show_labels:
                try:
                    if chart_type_str == 'pie':
                        ser_obj.data_labels.show_value = False
                        ser_obj.data_labels.show_percentage = True
                        try:
                            ser_obj.data_labels.number_format = '0%'
                        except Exception:
                            pass
                    else:
                        ser_obj.data_labels.show_value = True
                    lbl_color = (series_styles[si].get('data_label_color')
                                 if si < len(series_styles) else None)
                    if not lbl_color:
                        # Agent 6 renders labels outside the data mark, so default to a
                        # brand text color rather than white / series fill.
                        lbl_color = bt.get('body_color') or bt.get('primary_color', '#1A3C8F')
                    _lbl_size = min(max_chart_label_size, header_font_size) if allow_chart_fallback else max_chart_label_size
                    _lbl_pt = Pt(_lbl_size)
                    _lbl_rgb = hex_to_rgb(lbl_color)
                    # Set at series level first — this is the reliable path in python-pptx
                    try:
                        ser_obj.data_labels.font.size = _lbl_pt
                        ser_obj.data_labels.font.color.rgb = _lbl_rgb
                    except Exception:
                        pass
                    # Also apply per-label for charts where series-level is ignored
                    for lbl in ser_obj.data_labels:
                        try:
                            lbl.font.size = _lbl_pt
                            lbl.font.color.rgb = _lbl_rgb
                        except Exception:
                            pass
                    # Force data labels outside the data mark for all chart types where
                    # PowerPoint supports dLblPos so labels render in the slide background area.
                    if chart_type_str == 'horizontal_bar':
                        try:
                            dLbls_el = ser_obj.data_labels._element
                            pos_el = dLbls_el.find(nsmap.qn('c:dLblPos'))
                            if pos_el is None:
                                pos_el = etree.SubElement(dLbls_el, nsmap.qn('c:dLblPos'))
                            pos_el.set('val', 'outEnd')
                        except Exception:
                            pass
                    else:
                        try:
                            dLbls_el = ser_obj.data_labels._element
                            pos_el = dLbls_el.find(nsmap.qn('c:dLblPos'))
                            if pos_el is None:
                                pos_el = etree.SubElement(dLbls_el, nsmap.qn('c:dLblPos'))
                            pos_el.set('val', 'outEnd')
                        except Exception:
                            pass
                except Exception:
                    pass
    except Exception:
        pass

    # Pie: color each slice individually and rebuild dLbls via raw XML so
    # Outside End positioning is honored reliably by PowerPoint.
    if chart_type_str == 'pie':
        try:
            for ser_obj in chart.series:
                for pi, point_obj in enumerate(ser_obj.points):
                    c_hex = (series_styles[pi].get('fill_color')
                             if pi < len(series_styles) else None)
                    if not c_hex and allow_chart_fallback:
                        c_hex = chart_palette[pi % len(chart_palette)]
                    if c_hex:
                        point_obj.format.fill.solid()
                        point_obj.format.fill.fore_color.rgb = hex_to_rgb(c_hex)
        except Exception:
            pass

        if show_labels:
            try:
                pie_label_color = next(
                    (s.get('data_label_color') for s in series_styles if s.get('data_label_color')),
                    None
                ) or bt.get('body_color') or bt.get('primary_color') or '#1A1A1A'

                for si, ser_obj in enumerate(chart.series):
                    ser_el = ser_obj._element
                    values = []
                    if si < len(series_data):
                        values = [
                            float(v) if v is not None else 0.0
                            for v in (series_data[si].get('values') or [])
                        ]
                    while len(values) < len(categories):
                        values.append(0.0)
                    values = values[:len(categories)]

                    for existing in ser_el.findall(nsmap.qn('c:dLbls')):
                        ser_el.remove(existing)

                    dLbls_el = etree.Element(nsmap.qn('c:dLbls'))

                    for pi, val in enumerate(values):
                        if val == 0.0:
                            dLbl = etree.SubElement(dLbls_el, nsmap.qn('c:dLbl'))
                            etree.SubElement(dLbl, nsmap.qn('c:idx')).set('val', str(pi))
                            etree.SubElement(dLbl, nsmap.qn('c:delete')).set('val', '1')

                    numFmt = etree.SubElement(dLbls_el, nsmap.qn('c:numFmt'))
                    numFmt.set('formatCode', '0%')
                    numFmt.set('sourceLinked', '0')

                    txPr = etree.SubElement(dLbls_el, nsmap.qn('c:txPr'))
                    etree.SubElement(txPr, nsmap.qn('a:bodyPr'))
                    etree.SubElement(txPr, nsmap.qn('a:lstStyle'))
                    p_el = etree.SubElement(txPr, nsmap.qn('a:p'))
                    pPr_el = etree.SubElement(p_el, nsmap.qn('a:pPr'))
                    defRPr = etree.SubElement(pPr_el, nsmap.qn('a:defRPr'))
                    defRPr.set('sz', str(int(max(7, max_chart_label_size) * 100)))
                    solidFill = etree.SubElement(defRPr, nsmap.qn('a:solidFill'))
                    srgbClr = etree.SubElement(solidFill, nsmap.qn('a:srgbClr'))
                    srgbClr.set('val', pie_label_color.lstrip('#'))

                    etree.SubElement(dLbls_el, nsmap.qn('c:dLblPos')).set('val', 'outEnd')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showLdrLines')).set('val', '1')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showLegendKey')).set('val', '0')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showVal')).set('val', '0')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showCatName')).set('val', '0')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showSerName')).set('val', '0')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showPercent')).set('val', '1')
                    etree.SubElement(dLbls_el, nsmap.qn('c:showBubbleSize')).set('val', '0')

                    cat_el = ser_el.find(nsmap.qn('c:cat'))
                    if cat_el is not None:
                        ser_el.insert(list(ser_el).index(cat_el), dLbls_el)
                    else:
                        ser_el.append(dLbls_el)
            except Exception as e:
                print(f'render_chart pie data label XML error: {e}')

    # Axis font sizes + category label rotation for many categories
    try:
        ax_font_size = cs.get('axis_font_size', 9)
        if allow_chart_fallback:
            ax_font_size = min(ax_font_size, max_chart_label_size)
        ax_color     = cs.get('axis_color', '#000000')
        for ax in (chart.category_axis, chart.value_axis):
            try:
                ax.tick_labels.font.size      = Pt(ax_font_size)
                ax.tick_labels.font.color.rgb = hex_to_rgb(ax_color)
            except Exception:
                pass

        # Rotate category labels — prefer Agent 5 pre-computed rotation; fall back to heuristic.
        n_cats = len(categories)
        if _computed_cat_rotation is not None:
            rotation_deg = float(_computed_cat_rotation)
        elif allow_chart_fallback and chart_type_str in ('bar', 'line', 'clustered_bar') and n_cats > 6:
            rotation_deg = -45.0
        else:
            rotation_deg = None
        if rotation_deg is not None:
            try:
                # -2700000 EMU = -270 degrees = 45° slanted (OOXML uses 1/60000 degree units * -1)
                # rotation: -45 degrees = -2700000 in OOXML tickLblSkip/txPr
                rotation_emu = str(int(rotation_deg * 60000))
                cat_ax = chart.category_axis
                txPr = cat_ax.tick_labels._txPr
                if txPr is None:
                    # Create txPr if not present
                    txPr_elem = etree.SubElement(
                        cat_ax.tick_labels._element,
                        nsmap.qn('c:txPr')
                    )
                    etree.SubElement(txPr_elem, nsmap.qn('a:bodyPr')).set('rot', rotation_emu)
                    etree.SubElement(txPr_elem, nsmap.qn('a:lstStyle'))
                    p = etree.SubElement(txPr_elem, nsmap.qn('a:p'))
                    pr = etree.SubElement(p, nsmap.qn('a:pPr'))
                    etree.SubElement(pr, nsmap.qn('a:defRPr'))
                else:
                    bodyPr = txPr.find(nsmap.qn('a:bodyPr'))
                    if bodyPr is not None:
                        bodyPr.set('rot', rotation_emu)
            except Exception:
                pass
    except Exception:
        pass


def render_cards(slide, artifact, bt):
    """Render a cards artifact — each card as a styled rectangle + text."""
    cs       = artifact.get('card_style', {})
    ts       = artifact.get('title_style', {})
    sub_s    = artifact.get('subtitle_style', {})
    bs       = artifact.get('body_style', {})
    frames   = artifact.get('card_frames', [])
    cards    = artifact.get('cards', [])
    layout   = str(artifact.get('cards_layout', 'row') or 'row').lower()
    allow_card_fallback = renderer_fallback_allowed(artifact) or not frames

    def _num(v, default=None):
        try:
            if v is None:
                return default
            return float(v)
        except Exception:
            return default

    def _valid_rect(rect):
        return isinstance(rect, dict) and all(_num(rect.get(k)) is not None for k in ('x', 'y', 'w', 'h'))

    def _resolve_container_rect():
        container = artifact.get('container') or {}
        if _valid_rect(container):
            return {
                'x': _num(container.get('x'), 0.0),
                'y': _num(container.get('y'), 0.0),
                'w': max(0.2, _num(container.get('w'), 0.2)),
                'h': max(0.2, _num(container.get('h'), 0.2)),
            }

        if all(_num(artifact.get(k)) is not None for k in ('x', 'y', 'w', 'h')):
            return {
                'x': _num(artifact.get('x'), 0.0),
                'y': _num(artifact.get('y'), 0.0),
                'w': max(0.2, _num(artifact.get('w'), 0.2)),
                'h': max(0.2, _num(artifact.get('h'), 0.2)),
            }

        valid_frames = [f for f in frames if _valid_rect(f)]
        if valid_frames:
            min_x = min(_num(f.get('x'), 0.0) for f in valid_frames)
            min_y = min(_num(f.get('y'), 0.0) for f in valid_frames)
            max_x = max(_num(f.get('x'), 0.0) + max(0.01, _num(f.get('w'), 0.01)) for f in valid_frames)
            max_y = max(_num(f.get('y'), 0.0) + max(0.01, _num(f.get('h'), 0.01)) for f in valid_frames)
            return {
                'x': min_x,
                'y': min_y,
                'w': max(0.2, max_x - min_x),
                'h': max(0.2, max_y - min_y),
            }
        return None

    def _auto_card_frames(count, rect, layout_mode):
        if count <= 0 or not rect:
            return []

        gap = CARD_GAP
        rx, ry, rw, rh = rect['x'], rect['y'], rect['w'], rect['h']

        if layout_mode == 'column':
            each_h = max(0.2, (rh - gap * (count - 1)) / max(count, 1))
            return [
                {'x': rx, 'y': ry + i * (each_h + gap), 'w': rw, 'h': each_h}
                for i in range(count)
            ]

        if layout_mode == 'grid':
            cols = 2 if count > 1 else 1
            rows = int((count + cols - 1) / cols)
            cell_w = max(0.2, (rw - gap * (cols - 1)) / cols)
            cell_h = max(0.2, (rh - gap * (rows - 1)) / max(rows, 1))
            out = []
            for i in range(count):
                row = i // cols
                col = i % cols
                out.append({
                    'x': rx + col * (cell_w + gap),
                    'y': ry + row * (cell_h + gap),
                    'w': cell_w,
                    'h': cell_h,
                })
            return out

        each_w = max(0.2, (rw - gap * (count - 1)) / max(count, 1))
        return [
            {'x': rx + i * (each_w + gap), 'y': ry, 'w': each_w, 'h': rh}
            for i in range(count)
        ]

    container_rect = _resolve_container_rect()
    valid_frames = [f for f in frames if _valid_rect(f)]
    # If Agent 5 pre-computed card_frames (via computeArtifactInternals), use them directly.
    # Fall back to _auto_card_frames only when frames are absent or count mismatches.
    if len(valid_frames) == len(cards):
        frames = valid_frames   # use Agent 5 pre-computed frames as-is
    elif allow_card_fallback:
        frames = _auto_card_frames(len(cards), container_rect, layout)
        if not _valid_rect(artifact.get('container') or {}):
            artifact['container'] = container_rect
    else:
        frames = valid_frames

    fill_hex   = cs.get('fill_color', '#F5F5F5')
    border_hex = cs.get('border_color') or '#E0E0E0'
    border_w   = cs.get('border_width', 0.75)
    corner_r   = 0
    padding    = cs.get('internal_padding', CARD_INNER_PADDING)
    accent_w   = 0.07
    accent_gap = 0.08

    t_font  = ts.get('font_family', bt.get('title_font_family', 'Arial'))
    t_size  = ts.get('font_size', 10)
    t_color = ts.get('color', bt.get('primary_color', '#1A3C8F'))
    t_bold  = ts.get('font_weight', 'bold') in ('bold', 'semibold')

    su_font  = sub_s.get('font_family', bt.get('body_font_family', 'Arial'))
    su_size  = sub_s.get('font_size', 20)
    su_color = sub_s.get('color', '#000000')

    b_font  = bs.get('font_family', bt.get('body_font_family', 'Arial'))
    b_size  = bs.get('font_size', 9)
    b_color = bs.get('color', '#111111')

    # Sentinel fill colors for sentiment (override card_style.fill_color per card)
    _SENTIMENT_FILL = {
        'positive': '#EEF7F0',   # very light green
        'negative': '#FDF3F1',   # very light red
        'neutral':  None,        # use card_style fill as-is
    }
    _SENTIMENT_ACCENT = {
        'positive': bt.get('secondary_color', '#2D8A4E'),
        'negative': '#C0392B',
        'neutral':  bt.get('primary_color', '#1A3C8F'),
    }
    accent_palette = []
    for c in [bt.get('primary_color'), bt.get('secondary_color'), *(bt.get('accent_colors') or []), *(bt.get('chart_palette') or [])]:
        if c and c not in accent_palette:
            accent_palette.append(c)

    for fi, frame in enumerate(frames):
        fx = frame.get('x', 0)
        fy = frame.get('y', 0)
        fw = frame.get('w', 2)
        fh = frame.get('h', 1)

        card = cards[fi] if fi < len(cards) else {}

        # Derive per-card fill based on sentiment (if fill not explicitly set in card_style)
        sentiment  = card.get('sentiment', 'neutral')
        card_fill  = _SENTIMENT_FILL.get(sentiment) or fill_hex
        accent = (accent_palette[fi] if len(cards) > 1 and fi < len(accent_palette) else None) or _SENTIMENT_ACCENT.get(sentiment) or t_color

        # Card background
        add_filled_rect(slide, fx, fy, fw, fh,
                        fill_hex=card_fill,
                        border_hex=border_hex,
                        border_pt=border_w,
                        corner_radius=corner_r)

        # Left accent strip — brand-sequenced when multiple cards share an artifact
        add_filled_rect(slide, fx, fy, accent_w, fh, fill_hex=accent)

        # Proportional inner layout — works for any card height
        inner_left = fx + padding + accent_w + accent_gap
        inner_top  = fy + padding
        inner_w    = max(0.3, fw - padding * 2 - accent_w - accent_gap)
        inner_h    = fh - 2 * padding

        card_title = card.get('title', '')
        card_sub   = card.get('subtitle', '')
        card_body  = card.get('body', '')

        has_sub  = bool(card_sub)
        has_body = bool(card_body)

        # Short-card threshold: when card height is below 1.10" and there is a big
        # primary value (subtitle), a vertical stack makes the number crash into the
        # heading.  Switch to a horizontal split: primary value LEFT, heading+body RIGHT.
        USE_HORIZONTAL_SPLIT = inner_h < 1.0 and has_sub

        if inner_h > 0.05:
            if USE_HORIZONTAL_SPLIT:
                # ── Horizontal split layout ──────────────────────────────────────
                # Left column: primary value (big number)
                # Right column: title (heading) stacked above body
                split_gap   = 0.06
                left_col_w  = max(0.30, min(inner_w * 0.38, 0.65))
                right_col_w = max(0.20, inner_w - left_col_w - split_gap)
                right_col_x = inner_left + left_col_w + split_gap

                # Primary value — left column, vertically centred
                actual_su_size = estimate_fit_font_size(
                    str(card_sub), left_col_w, inner_h, su_size, 14
                ) if allow_card_fallback else su_size
                add_text_box(slide, inner_left, inner_top, left_col_w, inner_h,
                             str(card_sub), su_font, actual_su_size, True,
                             su_color, 'left', 'middle')

                # Title — top of right column
                title_h = max(0.16, inner_h * 0.36) if card_title else 0
                body_h  = max(0.10, inner_h - title_h - (TITLE_TO_SUBTITLE if card_title and card_body else 0)) if card_body else 0

                if card_title and title_h > 0:
                    actual_title_size = estimate_fit_font_size(
                        str(card_title), right_col_w, title_h, t_size, 7
                    ) if allow_card_fallback else t_size
                    add_text_box(slide, right_col_x, inner_top, right_col_w, title_h,
                                 str(card_title), t_font, actual_title_size, t_bold,
                                 ts.get('color', accent), 'left', 'top')

                body_text = str(card_body or '')
                if body_text and allow_card_fallback:
                    max_chars = max(30, int((right_col_w * 72 / 5.5) * 2))
                    if len(body_text) > max_chars:
                        body_text = body_text[:max_chars - 1].rstrip() + '…'
                if body_text and body_h > 0.05:
                    actual_body_size = estimate_fit_font_size(
                        body_text, right_col_w, body_h, b_size, 7
                    ) if allow_card_fallback else b_size
                    body_y = inner_top + title_h + (TITLE_TO_SUBTITLE if card_title and card_body else 0)
                    add_text_box(slide, right_col_x, body_y, right_col_w, body_h,
                                 body_text, b_font, actual_body_size, False,
                                 b_color, 'left', 'top')

            else:
                # ── Vertical stack layout (tall cards) ──────────────────────────
                title_h = max(0.18, inner_h * 0.16) if card_title else 0
                sub_h   = max(0.28, inner_h * 0.40) if card_sub else 0
                body_h  = max(0.12, inner_h - title_h - sub_h - (TITLE_TO_SUBTITLE if card_title and card_sub else 0) - (SUBTITLE_TO_BODY if card_sub and card_body else 0)) if card_body else 0

                actual_title_size = estimate_fit_font_size(str(card_title), max(0.3, fw - padding*2), max(0.14, title_h), t_size, 8) if (allow_card_fallback and card_title) else t_size
                actual_su_size = estimate_fit_font_size(str(card_sub), max(0.3, fw - padding*2), max(0.22, sub_h), su_size, 14) if (allow_card_fallback and sub_h > 0) else su_size

                body_text = str(card_body or '')
                if body_text and allow_card_fallback:
                    max_chars = max(30, int(((fw - padding * 2) * 72 / 5.5) * 2))
                    if len(body_text) > max_chars:
                        body_text = body_text[:max_chars - 1].rstrip() + '…'
                actual_body_size = estimate_fit_font_size(body_text, max(0.3, fw - padding*2), max(0.10, body_h), b_size, 7) if (allow_card_fallback and body_h > 0) else b_size

                title_y = inner_top
                if card_title and title_h > 0:
                    add_text_box(slide, inner_left, title_y, inner_w, title_h,
                                 str(card_title), t_font, actual_title_size, t_bold, ts.get('color', accent), 'left', 'top')

                if card_sub and sub_h > 0:
                    subtitle_y = inner_top + title_h + (TITLE_TO_SUBTITLE if card_title else 0)
                    add_text_box(slide, inner_left, subtitle_y, inner_w, sub_h,
                                 str(card_sub), su_font, actual_su_size, True, su_color, 'left', 'middle')

                if body_text and body_h > 0.05:
                    body_y = fy + fh - padding - body_h
                    add_text_box(slide, inner_left, body_y, inner_w, body_h,
                                 body_text, b_font, actual_body_size, False, b_color, 'left', 'bottom')


def render_workflow(slide, artifact, bt):
    """Render a workflow artifact — nodes as rectangles, connections as lines."""
    ws    = artifact.get('workflow_style', {})
    nodes = artifact.get('nodes', [])
    conns = artifact.get('connections', [])
    flow_direction = str(artifact.get('flow_direction', '') or '').lower()
    workflow_type = str(artifact.get('workflow_type', '') or '').lower()
    allow_workflow_fallback = renderer_fallback_allowed(artifact)

    def _num(v, default=0.0):
        try:
            return float(v)
        except Exception:
            return default

    def _has_rect(rect):
        return isinstance(rect, dict) and all(rect.get(k) is not None for k in ('x', 'y', 'w', 'h'))

    def _nodes_have_explicit_layout(nodes_):
        return bool(nodes_) and all(_has_rect(n) for n in nodes_)

    def _connections_have_explicit_paths(conns_):
        return all(isinstance(c.get('path'), list) and len(c.get('path', [])) >= 2 for c in conns_)

    def _workflow_bounds(nodes_, conns_):
        xs, ys, xe, ye = [], [], [], []
        for node in nodes_:
            nx = _num(node.get('x'))
            ny = _num(node.get('y'))
            nw = max(0.01, _num(node.get('w'), 2.0))
            nh = max(0.01, _num(node.get('h'), 0.8))
            xs.append(nx); ys.append(ny); xe.append(nx + nw); ye.append(ny + nh)
        for conn in conns_:
            for pt in conn.get('path', []):
                px = _num(pt.get('x'))
                py = _num(pt.get('y'))
                xs.append(px); ys.append(py); xe.append(px); ye.append(py)
        if not xs:
            return None
        return {
            'x': min(xs), 'y': min(ys),
            'w': max(0.01, max(xe) - min(xs)),
            'h': max(0.01, max(ye) - min(ys))
        }

    def _normalize_workflow_to_rect(nodes_, conns_, target_rect):
        source_rect = _workflow_bounds(nodes_, conns_)
        if not source_rect or not _has_rect(target_rect):
            return nodes_, conns_

        pad_x = min(0.12, max(0.03, target_rect['w'] * 0.04))
        pad_y = min(0.12, max(0.03, target_rect['h'] * 0.06))
        usable_w = max(0.05, target_rect['w'] - 2 * pad_x)
        usable_h = max(0.05, target_rect['h'] - 2 * pad_y)
        sx = usable_w / max(source_rect['w'], 0.01)
        sy = usable_h / max(source_rect['h'], 0.01)

        def _map_x(v):
            return target_rect['x'] + pad_x + (_num(v) - source_rect['x']) * sx

        def _map_y(v):
            return target_rect['y'] + pad_y + (_num(v) - source_rect['y']) * sy

        mapped_nodes = []
        for node in nodes_:
            mapped_nodes.append({
                **node,
                'x': _map_x(node.get('x')),
                'y': _map_y(node.get('y')),
                'w': max(0.20, _num(node.get('w'), 2.0) * sx),
                'h': max(0.20, _num(node.get('h'), 0.8) * sy),
            })

        mapped_conns = []
        for conn in conns_:
            mapped_conns.append({
                **conn,
                'path': [
                    {'x': _map_x(pt.get('x')), 'y': _map_y(pt.get('y'))}
                    for pt in conn.get('path', [])
                ]
            })

        return mapped_nodes, mapped_conns

    def _layout_horizontal_timeline(nodes_, target_rect):
        if not nodes_ or not _has_rect(target_rect):
            return nodes_, []

        count = max(1, len(nodes_))
        pad_x = min(0.10, max(0.04, target_rect['w'] * 0.025))
        pad_y = min(0.10, max(0.04, target_rect['h'] * 0.04))
        gap = min(0.12, max(0.05, target_rect['w'] * 0.02))
        usable_w = max(0.3, target_rect['w'] - 2 * pad_x - gap * (count - 1))
        node_w = max(0.75, usable_w / count)
        total_h = max(0.70, target_rect['h'] - 2 * pad_y)

        has_top_notes = any(bool(n.get('value', '')) for n in nodes_)
        has_bottom_notes = any(bool(n.get('description', '')) for n in nodes_)
        note_gap = 0.06
        top_h = max(0.16, min(0.34, total_h * 0.16)) if has_top_notes else 0.0
        desired_box_h = total_h * (0.34 if has_bottom_notes else 0.50)
        box_h = max(0.72, min(1.05, desired_box_h))
        bottom_h = max(0.0, total_h - box_h - top_h
                       - (note_gap if has_top_notes else 0.0)
                       - (note_gap if has_bottom_notes else 0.0))
        if has_bottom_notes and bottom_h < 0.18:
            deficit = 0.18 - bottom_h
            reducible = max(0.0, box_h - 0.62)
            take = min(deficit, reducible)
            box_h -= take
            bottom_h += take

        start_x = target_rect['x'] + pad_x
        node_y = target_rect['y'] + max(0.02, (target_rect['h'] - total_h) / 2.0)

        laid_out = []
        for i, node in enumerate(nodes_):
            box_y = node_y + top_h + (note_gap if has_top_notes else 0.0)
            laid_out.append({
                **node,
                'x': start_x + i * (node_w + gap),
                'y': node_y,
                'box_y': box_y,
                'box_h': box_h,
                'top_h': top_h,
                'bottom_h': bottom_h,
                'note_top_y': node_y,
                'note_top_h': top_h,
                'note_bottom_y': box_y + box_h + (note_gap if has_bottom_notes else 0.0),
                'note_bottom_h': bottom_h,
                'w': node_w,
                'h': total_h
            })

        conn_y = node_y + top_h + (note_gap if has_top_notes else 0.0) + box_h / 2.0
        laid_conns = []
        for i in range(len(laid_out) - 1):
            cur = laid_out[i]
            nxt = laid_out[i + 1]
            laid_conns.append({
                'from': cur.get('id'),
                'to': nxt.get('id'),
                'type': 'arrow',
                'path': [
                    {'x': cur['x'] + cur['w'], 'y': conn_y},
                    {'x': nxt['x'], 'y': conn_y}
                ]
            })

        return laid_out, laid_conns

    def _layout_vertical_flow(nodes_, target_rect, upward=False):
        if not nodes_ or not _has_rect(target_rect):
            return nodes_, []

        count = max(1, len(nodes_))
        pad_x = min(0.12, max(0.05, target_rect['w'] * 0.04))
        pad_y = min(0.12, max(0.05, target_rect['h'] * 0.04))
        gap_y = min(0.16, max(0.06, target_rect['h'] * 0.025))
        has_notes = any(bool(n.get('description', '') or n.get('value', '')) for n in nodes_)
        note_gap = 0.14 if has_notes else 0.0

        usable_w = max(0.4, target_rect['w'] - 2 * pad_x)
        box_w = min(max(1.05, usable_w * 0.40), max(1.05, usable_w - 0.75)) if has_notes else max(1.10, usable_w * 0.65)
        note_w = max(0.0, usable_w - box_w - note_gap)
        box_x = target_rect['x'] + pad_x
        note_x = box_x + box_w + note_gap

        band_h = max(0.80, (target_rect['h'] - 2 * pad_y - gap_y * (count - 1)) / count)
        box_h = max(0.54, min(0.85, band_h * 0.62))

        block_starts = [target_rect['y'] + pad_y + i * (band_h + gap_y) for i in range(count)]
        if upward:
            block_starts = list(reversed(block_starts))

        laid_out = []
        for i, node in enumerate(nodes_):
            block_y = block_starts[i]
            box_y = block_y + max(0.0, (band_h - box_h) / 2.0)
            laid_out.append({
                **node,
                'x': box_x,
                'y': box_y,
                'w': box_w,
                'h': box_h,
                'box_y': box_y,
                'box_h': box_h,
                'note_right_x': note_x,
                'note_right_y': block_y,
                'note_right_w': note_w,
                'note_right_h': band_h,
            })

        laid_conns = []
        for i in range(len(laid_out) - 1):
            cur = laid_out[i]
            nxt = laid_out[i + 1]
            laid_conns.append({
                'from': cur.get('id'),
                'to': nxt.get('id'),
                'type': 'arrow',
                'path': [
                    {'x': cur['x'] + cur['w'] / 2.0, 'y': cur['y'] + cur['h']},
                    {'x': nxt['x'] + nxt['w'] / 2.0, 'y': nxt['y']}
                ]
            })

        return laid_out, laid_conns

    target_rect = None
    container = artifact.get('container')
    if _has_rect(container):
        target_rect = {
            'x': _num(container.get('x')),
            'y': _num(container.get('y')),
            'w': max(0.05, _num(container.get('w'))),
            'h': max(0.05, _num(container.get('h')))
        }
    elif all(artifact.get(k) is not None for k in ('x', 'y', 'w', 'h')):
        target_rect = {
            'x': _num(artifact.get('x')),
            'y': _num(artifact.get('y')),
            'w': max(0.05, _num(artifact.get('w'))),
            'h': max(0.05, _num(artifact.get('h')))
        }
    needs_layout_fallback = (
        not _nodes_have_explicit_layout(nodes) or
        not _connections_have_explicit_paths(conns)
    )
    if target_rect and (allow_workflow_fallback or needs_layout_fallback):
        if flow_direction in ('left_to_right', 'horizontal') or workflow_type in ('timeline', 'roadmap'):
            nodes, conns = _layout_horizontal_timeline(nodes, target_rect)
        elif flow_direction in ('top_to_bottom', 'top_down', 'bottom_up'):
            nodes, conns = _layout_vertical_flow(nodes, target_rect, upward=(flow_direction == 'bottom_up'))
        else:
            nodes, conns = _normalize_workflow_to_rect(nodes, conns, target_rect)

    node_fill   = ws.get('node_fill_color', bt.get('primary_color', '#1A3C8F'))
    node_border = ws.get('node_border_color', '#FFFFFF')
    node_bw     = ws.get('node_border_width', 1.5)
    node_cr     = ws.get('node_corner_radius', 4)
    conn_color  = ws.get('connector_color', bt.get('primary_color', '#1A3C8F'))

    t_font  = ws.get('node_title_font_family', bt.get('title_font_family', 'Arial'))
    t_size  = ws.get('node_title_font_size', 11)
    t_color = ws.get('node_title_color', '#FFFFFF')
    t_bold  = ws.get('node_title_font_weight', 'bold') in ('bold', 'semibold')

    v_font  = ws.get('node_value_font_family', bt.get('body_font_family', 'Arial'))
    v_size  = ws.get('node_value_font_size', 10)
    v_color = ws.get('node_value_color', bt.get('body_color', '#000000'))
    top_note_color = bt.get('primary_color') or bt.get('title_color') or t_color

    def _fit_or_spec(text, width_in, height_in, spec_size, min_size):
        if allow_workflow_fallback:
            return estimate_fit_font_size(text, width_in, height_in, spec_size, min_size)
        return spec_size

    # Draw connections first (behind nodes)
    for conn in conns:
        path = conn.get('path', [])
        if len(path) >= 2:
            try:
                for i in range(len(path) - 1):
                    x1 = float(path[i].get('x', 0))
                    y1 = float(path[i].get('y', 0))
                    x2 = float(path[i + 1].get('x', 0))
                    y2 = float(path[i + 1].get('y', 0))
                    connector = slide.shapes.add_connector(
                        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
                        inches(x1), inches(y1), inches(x2), inches(y2)
                    )
                    connector.line.color.rgb = hex_to_rgb(conn_color)
                    connector.line.width     = Pt(ws.get('connector_width', 2))
            except Exception:
                pass

    # Draw nodes
    # Layout pattern:
    # - horizontal flows: optional note above, primary label inside box, optional note below
    # - vertical flows: primary label inside box, optional note to the right
    # - all other flows: primary/value inside box, description below
    body_color = bt.get('body_color', '#000000')
    is_horizontal_flow = flow_direction in ('left_to_right', 'horizontal') or workflow_type in ('timeline', 'roadmap')
    is_vertical_flow = flow_direction in ('top_to_bottom', 'top_down', 'bottom_up')
    effective_node_fill = (bt.get('primary_color') or bt.get('title_color') or node_fill) if is_horizontal_flow else node_fill
    for node in nodes:
        nx   = node.get('x', 0)
        ny   = node.get('y', 0)
        nw   = node.get('w', 2)
        nh   = node.get('h', 0.8)
        box_y = node.get('box_y', ny)
        bxh  = node.get('box_h', nh)

        # ── Colored fill box (label + value only) ────────────────────────────
        add_filled_rect(slide, nx, box_y, nw, bxh,
                        fill_hex=effective_node_fill,
                        border_hex=node_border,
                        border_pt=node_bw,
                        corner_radius=node_cr)

        # Node label (primary text — inside the colored box)
        label = node.get('label', node.get('id', ''))
        value = node.get('value', '')
        desc = node.get('description', '')
        has_value = bool(value)
        label_h = bxh * (0.55 if has_value else 0.70)
        val_h   = max(0.14, bxh - label_h - 0.10) if has_value else 0

        if label:
            label_size = _fit_or_spec(str(label), max(0.3, nw - 0.2), max(0.22, bxh - 0.12), t_size, 8)
            add_text_box(slide, nx + 0.1, box_y + 0.06, nw - 0.2, max(0.22, bxh - 0.12),
                         str(label), t_font, label_size, t_bold, t_color, 'center', 'middle')

        # Node value (secondary metric — also inside the colored box, below label)
        if is_horizontal_flow:
            top_h = max(0.0, float(node.get('note_top_h', 0.0) or 0.0))
            if value and top_h > 0.08:
                top_y = float(node.get('note_top_y', ny) or ny)
                value_size = _fit_or_spec(str(value), max(0.28, nw - 0.10), top_h, max(8, v_size), 7)
                add_text_box(slide, nx + 0.05, top_y, nw - 0.10, top_h,
                             str(value), v_font, value_size, True, top_note_color, 'center', 'middle')
        elif is_vertical_flow:
            note_x = float(node.get('note_right_x', nx + nw + 0.10) or (nx + nw + 0.10))
            note_y = float(node.get('note_right_y', ny) or ny)
            note_w = max(0.0, float(node.get('note_right_w', 0.0) or 0.0))
            note_h = max(0.0, float(node.get('note_right_h', nh) or nh))
            note_parts = [str(x).strip() for x in (value, desc) if str(x or '').strip()]
            note_text = '\n'.join(note_parts[:2])
            if note_text and note_w > 0.12 and note_h > 0.12:
                note_size = _fit_or_spec(note_text, max(0.24, note_w), note_h,
                                         max(7, v_size - 1), 6)
                add_text_box(slide, note_x, note_y, note_w, note_h,
                             note_text, v_font, note_size, False, body_color, 'left', 'middle')
        elif value:
            value_size = _fit_or_spec(str(value), max(0.3, nw - 0.2), max(0.14, val_h), max(7, v_size - 1), 7)
            add_text_box(slide, nx + 0.1, box_y + 0.06 + label_h, nw - 0.2, val_h,
                         str(value), v_font, value_size, False, t_color, 'center', 'middle')

        # ── Description text BELOW the colored box ───────────────────────────
        # Rendered outside the fill in body color so it reads as supporting detail,
        # not as part of the highlighted header block.
        if is_horizontal_flow:
            bottom_h = max(0.0, float(node.get('note_bottom_h', 0.0) or 0.0))
            if desc and bottom_h > 0.12:
                bottom_y = float(node.get('note_bottom_y', box_y + bxh) or (box_y + bxh))
                desc_size = _fit_or_spec(str(desc), max(0.28, nw - 0.12), bottom_h,
                                         max(7, v_size - 1), 6)
                add_text_box(slide, nx + 0.06, bottom_y, nw - 0.12, bottom_h,
                             str(desc), v_font, desc_size, False, body_color, 'center', 'top')
        elif not is_vertical_flow:
            desc_gap = 0.06
            desc_y   = box_y + bxh + desc_gap
            desc_h   = max(0.0, ny + nh - desc_y)
            if desc and desc_h > 0.12:
                desc_size = _fit_or_spec(str(desc), max(0.28, nw - 0.12), desc_h,
                                         max(7, v_size - 1), 6)
                add_text_box(slide, nx + 0.06, desc_y, nw - 0.12, desc_h,
                             str(desc), v_font, desc_size, False, body_color, 'center', 'top')


def render_table(slide, artifact, bt):
    """Render a table artifact using python-pptx native table."""
    x = artifact.get('x', 0)
    y = artifact.get('y', 0)
    w = artifact.get('w', 6)
    h = artifact.get('h', 2)

    ts      = artifact.get('table_style', {})
    headers = artifact.get('headers', [])
    rows    = artifact.get('rows', [])
    col_ws  = artifact.get('column_widths', [])
    row_hs  = artifact.get('row_heights', [])
    hl_rows = artifact.get('highlight_rows', [])
    pre_col_xs = artifact.get('column_x_positions', [])
    pre_row_ys = artifact.get('row_y_positions', [])
    header_cell_frames = artifact.get('header_cell_frames', [])
    body_cell_frames = artifact.get('body_cell_frames', [])

    if not headers:
        add_text_box(slide, x, y, w, h, 'Table (no headers)', 'Arial', 10,
                     False, '#888888', 'center', 'middle')
        return

    n_cols  = len(headers)
    n_rows  = len(rows) + 1   # +1 for header row
    data_rows = rows

    def _norm_sizes(values, count, total):
        nums = []
        for i in range(count):
            try:
                v = float(values[i]) if i < len(values) else 0
            except Exception:
                v = 0
            nums.append(max(0, v))
        summed = sum(nums)
        if summed > 0:
            return [(v / summed) * total for v in nums]
        return [total / max(count, 1)] * count

    def _auto_col_widths():
        weights = []
        for ci, hdr in enumerate(headers):
            max_len = len(str(hdr or ''))
            for row in data_rows:
                if ci < len(row):
                    max_len = max(max_len, len(str(row[ci] or '')))
            weights.append(max(6, min(max_len, 40)))
        total_weight = sum(weights) or 1
        return [w * (wt / total_weight) for wt in weights]

    # Use Agent 5 pre-computed column widths when available; else auto-compute.
    pre_col_ws      = artifact.get('column_widths', [])
    pre_col_types   = artifact.get('column_types', [])
    pre_col_aligns  = artifact.get('column_alignments', [])
    pre_row_hs      = artifact.get('row_heights', [])
    pre_hdr_h       = artifact.get('header_row_height', None)

    if header_cell_frames and len(header_cell_frames) == n_cols:
        norm_col_ws = [float(cf.get('w', 0)) for cf in header_cell_frames]
    elif pre_col_ws and len(pre_col_ws) == n_cols:
        norm_col_ws = list(pre_col_ws)
    else:
        norm_col_ws = _auto_col_widths()

    if header_cell_frames and len(header_cell_frames) == n_cols:
        x = float(header_cell_frames[0].get('x', x))
        y = float(header_cell_frames[0].get('y', y))
        if pre_col_xs and len(pre_col_xs) == n_cols:
            x = float(pre_col_xs[0])

    if pre_hdr_h is not None and pre_row_hs and len(pre_row_hs) == len(data_rows):
        # Build full row heights list: [header_h, data_row_h...]
        norm_row_hs = [float(pre_hdr_h)] + [float(rh) for rh in pre_row_hs]
        row_total = sum(norm_row_hs)
        if row_total > h:
            norm_row_hs = [rh * (h / row_total) for rh in norm_row_hs]
    elif header_cell_frames and body_cell_frames:
        header_h = float(header_cell_frames[0].get('h', pre_hdr_h or 0.35))
        body_hs = []
        for row_frames in body_cell_frames:
            if row_frames:
                body_hs.append(float(row_frames[0].get('h', 0)))
        norm_row_hs = [header_h] + body_hs
    else:
        norm_row_hs = [max(TABLE_MIN_ROW_HEIGHT, rh) for rh in _norm_sizes(row_hs, n_rows, h)]
        row_total = sum(norm_row_hs)
        if row_total > h:
            norm_row_hs = [rh * (h / row_total) for rh in norm_row_hs]

    if pre_row_ys and len(pre_row_ys) >= 1:
        y = float(pre_row_ys[0])

    w = sum(norm_col_ws) if norm_col_ws else w
    h = sum(norm_row_hs) if norm_row_hs else h

    try:
        table_shape = slide.shapes.add_table(n_rows, n_cols, inches(x), inches(y), inches(w), inches(h))
        table       = table_shape.table

        # Column widths
        for ci, cw in enumerate(norm_col_ws):
            table.columns[ci].width = inches(cw)

        # Row heights
        for ri, rh in enumerate(norm_row_hs):
            table.rows[ri].height = inches(rh)

        # Header row styles
        h_fill  = ts.get('header_fill_color', bt.get('primary_color', '#1A3C8F'))
        h_text  = ts.get('header_text_color', '#FFFFFF')
        h_font  = ts.get('header_font_family', bt.get('title_font_family', 'Arial'))
        h_size  = ts.get('header_font_size', 11)

        b_fill  = ts.get('body_fill_color', '#FFFFFF')
        b_alt   = ts.get('body_alt_fill_color')
        b_text  = ts.get('body_text_color', '#111111')
        b_font  = ts.get('body_font_family', bt.get('body_font_family', 'Arial'))
        b_size  = ts.get('body_font_size', 10)
        hl_fill = ts.get('highlight_fill_color')
        grid_c  = ts.get('grid_color', '#CCCCCC')
        cell_p  = float(ts.get('cell_padding', CELL_PADDING) or CELL_PADDING)

        # Detect which columns are numeric (right-align) vs text (left-align).
        # Use Agent 5 pre-computed column_types/alignments when available.
        if pre_col_types and len(pre_col_types) == n_cols:
            col_is_numeric = [t == 'numeric' for t in pre_col_types]
        else:
            def _is_numeric_col(col_idx):
                import re
                _num_pat = re.compile(r'^[\s₹$€£¥\-\+]?[\d,\.]+[%KMBCr\s]*$')
                hits = sum(1 for r in data_rows
                           if col_idx < len(r) and _num_pat.match(str(r[col_idx]).strip()))
                return (hits / max(len(data_rows), 1)) > 0.5

            col_is_numeric = [_is_numeric_col(ci) for ci in range(n_cols)]
            # First column is always text (entity/label column)
            if n_cols > 0:
                col_is_numeric[0] = False

        cell_pad_emu = int(Emu(cell_p * 914400))  # cell_p in inches → EMU

        def _apply_cell_margin(cell, pad_emu):
            """Apply uniform cell margin via OOXML tcPr."""
            try:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for attr in ('marL', 'marR', 'marT', 'marB'):
                    tcPr.set(attr, str(pad_emu))
            except Exception:
                pass

        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = str(hdr)
            enable_text_fit(cell.text_frame)
            try:
                cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(h_fill)
            _apply_cell_margin(cell, cell_pad_emu)
            try:
                run = cell.text_frame.paragraphs[0].runs[0]
                fit_size = estimate_fit_font_size(
                    hdr,
                    norm_col_ws[ci],
                    norm_row_hs[0],
                    h_size,
                    8
                )
                set_font(run, h_font, fit_size, True, False, h_text)
                # Header row: center-align all columns
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            except Exception:
                pass

        # Body rows
        for ri, row in enumerate(data_rows):
            row_idx = ri + 1
            use_alt = (b_alt and ri % 2 == 1)
            use_hl  = (hl_fill and ri in hl_rows)
            row_fill = hl_fill if use_hl else (b_alt if use_alt else b_fill)

            for ci in range(n_cols):
                cell = table.cell(row_idx, ci)
                cell_text = str(row[ci]) if ci < len(row) else ''
                cell.text = cell_text
                enable_text_fit(cell.text_frame)
                try:
                    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                except Exception:
                    pass
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(row_fill or b_fill)
                _apply_cell_margin(cell, cell_pad_emu)
                try:
                    run = cell.text_frame.paragraphs[0].runs[0]
                    is_highlight = ri in hl_rows
                    font_size = b_size + 0.5 if is_highlight else b_size
                    fit_size = estimate_fit_font_size(
                        cell_text,
                        norm_col_ws[ci],
                        norm_row_hs[row_idx],
                        font_size,
                        7
                    )
                    text_color = ts.get('highlight_text_color', b_text) if is_highlight else b_text
                    set_font(run, b_font, fit_size, is_highlight, False, text_color)
                    # Numeric columns → right-align; text columns → left-align.
                    # Prefer Agent 5 pre-computed alignments when available.
                    if pre_col_aligns and ci < len(pre_col_aligns):
                        if pre_col_aligns[ci] == 'right':
                            align = PP_ALIGN.RIGHT
                        elif pre_col_aligns[ci] == 'center':
                            align = PP_ALIGN.CENTER
                        else:
                            align = PP_ALIGN.LEFT
                    else:
                        align = PP_ALIGN.RIGHT if col_is_numeric[ci] else PP_ALIGN.LEFT
                    cell.text_frame.paragraphs[0].alignment = align
                except Exception:
                    pass

    except Exception as e:
        # Fallback — plain text box
        all_text = ' | '.join(headers) + '\n' + '\n'.join([' | '.join([str(c) for c in r]) for r in rows])
        add_text_box(slide, x, y, w, h, all_text, 'Arial', 9, False, '#111111')


def _write_heading_to_header_ph(slide, heading_text, header_ph_idx, bt, header_style='underline'):
    """Write heading text into the layout's paired header placeholder."""
    if not heading_text or header_ph_idx is None:
        return None
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == header_ph_idx:
                ph_x = ph.left / 914400
                ph_y = ph.top / 914400
                ph_w = ph.width / 914400
                ph_h = ph.height / 914400
                ph.text_frame.text = str(heading_text)
                try:
                    run = ph.text_frame.paragraphs[0].runs[0]
                    from pptx.util import Pt
                    from pptx.dml.color import RGBColor
                    run.font.name = bt.get('title_font_family', 'Arial')
                    run.font.size = Pt(10)
                    run.font.bold = True
                    color_hex = bt.get('primary_color', '#1A3C8F').lstrip('#')
                    run.font.color.rgb = RGBColor(
                        int(color_hex[0:2], 16),
                        int(color_hex[2:4], 16),
                        int(color_hex[4:6], 16)
                    )
                except Exception:
                    pass
                try:
                    font_size = 10
                    text_h = estimate_header_block_height(heading_text, ph_w, font_size)
                    if header_style == 'brand_fill':
                        fill_color = bt.get('primary_color', '#1A3C8F')
                        text_color = '#FFFFFF' if is_dark_color(fill_color) else '#111111'
                        # Hide placeholder text and render an explicit fill header above it.
                        try:
                            ph.text_frame.clear()
                        except Exception:
                            pass
                        add_filled_rect(slide, ph_x, ph_y, ph_w, HEADER_HEIGHT, fill_hex=fill_color)
                        add_text_box(
                            slide,
                            ph_x + MIN_TEXT_MARGIN,
                            ph_y,
                            max(0.2, ph_w - MIN_TEXT_MARGIN * 2),
                            HEADER_HEIGHT,
                            str(heading_text),
                            bt.get('title_font_family', 'Arial'),
                            font_size,
                            True,
                            text_color,
                            'left',
                            'middle'
                        )
                        return ph_y + HEADER_HEIGHT
                    else:
                        rule_y = ph_y + text_h + 0.02
                        add_filled_rect(
                            slide,
                            ph_x,
                            rule_y,
                            ph_w,
                            0.03,
                            fill_hex=bt.get('primary_color', '#1A3C8F')
                        )
                        return rule_y + 0.03
                except Exception:
                    pass
                return ph_y + max(text_h, HEADER_HEIGHT)
    except Exception as e:
        print(f'_write_heading_to_header_ph error (idx={header_ph_idx}):', e)
    return None


def render_artifact(slide, artifact, bt, ph_frame=None, header_ph_idx=None, header_style='underline', slide_w=13.33, slide_h=7.5):
    """
    Dispatch to the correct renderer based on artifact type.

    ph_frame      — dict with x/y/w/h keys (inches) pre-read from the layout placeholder
                    BEFORE that placeholder was removed from the slide.  When set, these
                    bounds override the artifact's own x/y/w/h (e.g. when Agent 5 left
                    them null in layout mode, trusting the placeholder for positioning).
    header_ph_idx — layout mode: write artifact heading into the paired header placeholder
                    that was preserved in the slide (not removed).
    """
    t = (artifact.get('type') or '').lower()
    artifact_header_style = infer_artifact_header_style(t)

    artifact = dict(artifact)   # shallow copy — don't mutate the spec

    # Apply pre-saved placeholder bounds when the artifact's own coords are null/missing
    if ph_frame is not None:
        artifact['x'] = ph_frame['x']
        artifact['y'] = ph_frame['y']
        artifact['w'] = ph_frame['w']
        artifact['h'] = ph_frame['h']

    # Route artifact heading into the layout's paired header placeholder when available
    heading_handled = False
    placeholder_header_bottom = None
    inline_header_rendered = False
    rendered_header_bottom = None
    header_block = artifact.get('header_block') or {}
    use_placeholder_header = bool(header_block.get('placeholder_ref'))
    if header_ph_idx is not None and t != 'cards' and use_placeholder_header:
        if t == 'insight_text':
            heading_text = artifact.get('heading') or artifact.get('insight_header', '')
        elif t == 'chart':
            heading_text = artifact.get('chart_header', '')
        elif t == 'table':
            heading_text = artifact.get('table_header', '')
        elif t == 'workflow':
            heading_text = artifact.get('workflow_header', '')
        else:
            heading_text = ''
        if heading_text:
            placeholder_header_bottom = _write_heading_to_header_ph(slide, heading_text, header_ph_idx, bt, header_style=artifact_header_style)
            heading_handled = placeholder_header_bottom is not None

    # Render header_block only when the heading wasn't routed to a layout placeholder.
    # NOTE: if placeholder_ref=True but heading_handled=False it means no real placeholder
    # was available — fall back to rendering the header inline ABOVE the artifact so it
    # sits outside any filled/bordered box rather than inside it.
    if t != 'cards' and not heading_handled:
        hb = dict(header_block)
        # Align header x/w to the artifact's effective bounds.
        # For workflow, bounds live in 'container', not top-level x/y/w/h — use container
        # as the authoritative source so the header never inherits a full-slide ph_frame width.
        if hb:
            _eff_x = artifact.get('x')
            _eff_y = artifact.get('y')
            _eff_w = artifact.get('w')
            _eff_h = artifact.get('h')
            if t == 'workflow' and (_eff_x is None or _eff_w is None):
                _ctr = artifact.get('container') or {}
                _eff_x = _ctr.get('x', _eff_x)
                _eff_y = _ctr.get('y', _eff_y)
                _eff_w = _ctr.get('w', _eff_w)
                _eff_h = _ctr.get('h', _eff_h)
            if _eff_x is not None and _eff_w is not None:
                hb['x'] = _eff_x
                hb['w'] = _eff_w
            if _eff_y is not None and hb.get('y') is None:
                hb['y'] = _eff_y
            if _eff_h is not None and hb.get('h') is None:
                hb['h'] = min(HEADER_HEIGHT, max(0.2, float(_eff_h)))
        if hb:
            # Always render inline when heading_handled=False — even when placeholder_ref=True,
            # since we only reach here when no placeholder actually received the text.
            # Strip placeholder_ref so render_header_block treats this as a normal inline header.
            hb['placeholder_ref'] = False
            rendered_header_bottom = render_header_block(slide, hb, bt, header_style=artifact_header_style)
            inline_header_rendered = True

    suppress_internal_heading = heading_handled or inline_header_rendered
    if inline_header_rendered and header_block:
        gap = 0.08
        header_bottom = float(rendered_header_bottom or 0) + gap
        art_y = float(artifact.get('y', 0) or 0)
        art_h = float(artifact.get('h', 0) or 0)
        if art_h > 0 and header_bottom > art_y:
            delta = header_bottom - art_y
            artifact['y'] = header_bottom
            artifact['h'] = max(0.2, art_h - delta)
    elif heading_handled:
        # Only introduce the minimum gap required after the actual rendered header.
        art_y = float(artifact.get('y', 0) or 0)
        art_h = float(artifact.get('h', 0) or 0)
        required_top = float(placeholder_header_bottom or art_y) + HEADER_TO_ARTIFACT
        if art_h > 0 and required_top > art_y:
            delta = required_top - art_y
            artifact['y'] = required_top
            artifact['h'] = max(0.2, art_h - delta)

    try:
        if   t == 'insight_text': render_insight_text(slide, artifact, bt,
                                                       suppress_heading=suppress_internal_heading)
        elif t == 'chart':        render_chart(slide, artifact, bt,
                                               suppress_heading=suppress_internal_heading,
                                               slide_w=slide_w, slide_h=slide_h)
        elif t == 'cards':        render_cards(slide, artifact, bt)
        elif t == 'workflow':     render_workflow(slide, artifact, bt)
        elif t == 'table':        render_table(slide, artifact, bt)
    except Exception as e:
        print(f'render_artifact error ({t}):', e)
        traceback.print_exc()


# ─── SLIDE BUILDER ─────────────────────────────────────────────────────────────

def _write_speaker_note(slide, note_text):
    if note_text:
        try:
            slide.notes_slide.notes_text_frame.text = str(note_text)
        except Exception:
            pass


def _remove_content_placeholders(slide, keep_indices=None):
    """
    Upfront removal of non-system content/body placeholders from the slide XML.

    Layout templates often have descriptive default text baked into their
    content placeholders (e.g. "z1 · summary", "primary") for design guidance.
    Setting text = '' is unreliable because python-pptx may fall through to the
    layout's inherited txBody.  The only guaranteed fix is to remove the
    placeholder shape element from the slide's spTree entirely.

    keep_indices — set of placeholder.idx values to PRESERVE (these will be
                   written to by _write_heading_to_header_ph for heading text).
                   All other non-system placeholders are removed.
    """
    from pptx.enum.shapes import PP_PLACEHOLDER as _PH
    _SYSTEM_TYPES = {_PH.TITLE, _PH.CENTER_TITLE, _PH.SUBTITLE,
                     _PH.DATE, _PH.FOOTER, _PH.SLIDE_NUMBER}
    keep_indices = keep_indices or set()

    sp_tree = slide.shapes._spTree
    to_remove = []
    for ph in list(slide.placeholders):
        try:
            ph_type = ph.placeholder_format.type
            ph_idx  = ph.placeholder_format.idx
            if ph_type not in _SYSTEM_TYPES and ph_idx not in keep_indices:
                to_remove.append(ph._element)
        except Exception:
            pass
    for elem in to_remove:
        try:
            sp_tree.remove(elem)
        except Exception:
            pass


def _remove_empty_placeholders(slide):
    """
    Remove placeholder shapes that carry no content.
    After rendering, any layout placeholder that was not explicitly written to
    would appear as a 'Click to add text' ghost box overlaid on the slide.
    This is safe because we render all content as free shapes or via
    place_in_placeholder — so a placeholder is only needed when it has text.
    """
    sp_tree = slide.shapes._spTree
    to_remove = []
    for ph in list(slide.placeholders):
        try:
            if not ph.text_frame.text.strip():
                to_remove.append(ph._element)
        except Exception:
            # Non-text placeholder (picture / object type) — treat as unfilled
            to_remove.append(ph._element)
    for elem in to_remove:
        try:
            sp_tree.remove(elem)
        except Exception:
            pass


def _sanitise_system_placeholders(slide, slide_number):
    """
    After add_slide(), fix placeholders that carry stale template text:
      - dt  (date/time)   → blank — presentation-level date field or master handles it
      - ftr (footer)      → blank — confidentiality text is on the master; don't duplicate
      - sldNum (slide #)  → write actual slide number so it shows correctly in PDF/export
    """
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER as _PH
        for ph in slide.placeholders:
            ph_type = ph.placeholder_format.type
            if ph_type in (_PH.DATE, _PH.FOOTER):
                try:
                    tf = ph.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    # Also clear via XML to strip any fld elements carrying old date
                    for txBody in ph._element.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}txBody'):
                        for p_el in list(txBody):
                            tag = p_el.tag.split('}')[-1]
                            if tag == 'p':
                                for child in list(p_el):
                                    ctag = child.tag.split('}')[-1]
                                    if ctag in ('r', 'fld'):
                                        p_el.remove(child)
                except Exception:
                    pass
            elif ph_type == _PH.SLIDE_NUMBER:
                try:
                    ph.text_frame.text = str(slide_number)
                except Exception:
                    pass
    except Exception:
        pass


def normalize_zone_artifact_stack(zone):
    """For multi-artifact zones, enforce a stacked 65/35 split inside the zone frame."""
    if not isinstance(zone, dict):
        return zone
    artifacts = [dict(a) for a in (zone.get('artifacts') or [])]
    frame = zone.get('frame') or {}
    padding = frame.get('padding') or {}
    if len(artifacts) < 2 or not all(k in frame for k in ('x', 'y', 'w', 'h')):
        return { **zone, 'artifacts': artifacts }

    left = float(frame.get('x', 0)) + float(padding.get('left', INTERNAL_PADDING) or 0)
    top = float(frame.get('y', 0)) + float(padding.get('top', ZONE_TOP_OFFSET) or 0)
    width = max(0.2, float(frame.get('w', 0)) - float(padding.get('left', INTERNAL_PADDING) or 0) - float(padding.get('right', INTERNAL_PADDING) or 0))
    height = max(0.2, float(frame.get('h', 0)) - float(padding.get('top', ZONE_TOP_OFFSET) or 0) - float(padding.get('bottom', INTERNAL_PADDING) or 0))

    gap = ARTIFACT_TO_ARTIFACT
    primary_h = max(0.2, (height - gap) * 0.65)
    secondary_band_h = max(0.2, height - primary_h - gap)

    laid_out = []
    for idx, art in enumerate(artifacts):
        # Agent 5 pre-computed this artifact's bounds — skip override
        if all(art.get(k) is not None for k in ('x', 'y', 'w', 'h')):
            laid_out.append(art)
            continue
        if idx == 0:
            ax, ay, aw, ah = left, top, width, primary_h
        else:
            remaining = len(artifacts) - 1
            each_h = max(0.15, (secondary_band_h - max(0, remaining - 1) * gap) / max(remaining, 1))
            ax = left
            ay = top + primary_h + gap + (idx - 1) * (each_h + gap)
            aw = width
            ah = each_h
        art['x'] = ax
        art['y'] = ay
        art['w'] = aw
        art['h'] = ah
        if art.get('type') == 'workflow':
            art['container'] = { 'x': ax, 'y': ay, 'w': aw, 'h': ah }
        laid_out.append(art)

    return { **zone, 'artifacts': laid_out }


def _shift_blocks_for_title_gap(slide, blocks, use_template):
    """
    In template mode, Agent 5 estimates title_block.h from limited info.
    The actual title placeholder in the template has a fixed height from the
    layout XML — always deterministic, no font estimation needed.

    Algorithm (per user spec):
      1. Place title/subtitle (done by caller before render_blocks).
         We read placeholder bounds BEFORE text is placed — bounds are fixed
         by the layout and don't change when text is written.
      2. Compute actual_header_bottom = bottom of title (or subtitle if present)
         placeholder, in inches.
      3. Compare with content_start_y = min Y of all non-header blocks.
      4. If gap = content_start_y - actual_header_bottom < MIN_GAP_IN:
         shift every non-header block down by (MIN_GAP_IN - gap).

    MIN_GAP_IN ≈ 2pt ≈ 0.028" — just enough to guarantee no visual overlap.
    No font size, no line-count estimation needed.
    """
    if not use_template:
        return blocks

    title_block = next(
        (b for b in blocks if b.get('block_type') == 'title' and b.get('text')), None
    )
    if not title_block:
        return blocks

    # Finalized contract uses 2px, not 2pt. Assume standard Office/render DPI.
    MIN_GAP_IN = 32 / 96.0

    # Use the block's own y+h as the header bottom — these are Agent 5's pre-computed
    # coordinates and are the most direct measure of where the header ends.
    # No placeholder measurement or font estimation needed.
    subtitle_block = next(
        (b for b in blocks if b.get('block_type') == 'subtitle' and b.get('text')),
        None
    )

    try:
        header_bottom = float(title_block['y']) + float(title_block['h'])
    except (KeyError, TypeError, ValueError):
        return blocks

    if subtitle_block:
        try:
            sub_bottom = float(subtitle_block['y']) + float(subtitle_block['h'])
            header_bottom = max(header_bottom, sub_bottom)
        except (KeyError, TypeError, ValueError):
            pass

    # Earliest Y of all non-header content blocks
    content_blocks = [
        b for b in blocks
        if b.get('block_type') not in ('title', 'subtitle')
        and b.get('y') is not None
    ]
    if not content_blocks:
        return blocks

    content_start_y = min(float(b['y']) for b in content_blocks)
    gap   = content_start_y - header_bottom
    shift = round(MIN_GAP_IN - gap, 16)   # positive → content too close / overlapping

    if shift <= 0:
        return blocks   # content already clears the header by at least MIN_GAP_IN

    print(f'[title gap fix] header_bottom={header_bottom:.3f}" '
          f'content_start={content_start_y:.3f}" gap={gap:.3f}" shift={shift:.3f}"')

    shifted = []
    for b in blocks:
        if b.get('block_type') in ('title', 'subtitle'):
            shifted.append(b)
        elif b.get('y') is not None:
            shifted.append({**b, 'y': round(float(b['y']) + shift, 3)})
        else:
            shifted.append(b)
    return shifted


def _ensure_subtitle_placeholder_clears_title(slide, blocks):
    """
    In template mode, a long wrapped title can visually extend into the subtitle
    placeholder area. Before placing subtitle text, nudge the subtitle
    placeholder down so its top clears the estimated rendered title text bottom.
    """
    EMU = 914400.0
    MIN_GAP_IN = 2 / 96.0

    title_block = next(
        (b for b in blocks if b.get('block_type') == 'title' and b.get('text')),
        None
    )
    subtitle_block = next(
        (b for b in blocks if b.get('block_type') == 'subtitle' and b.get('text')),
        None
    )
    if not title_block or not subtitle_block:
        return

    def _ph(idx):
        try:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == idx:
                    return ph
        except Exception:
            pass
        return None

    title_ph = _ph(0)
    subtitle_ph = _ph(1)
    if not title_ph or not subtitle_ph:
        return

    try:
        title_top = title_ph.top / EMU
        title_w = max(0.5, title_ph.width / EMU)
        font_size = int(title_block.get('font_size') or 18)
        lines = estimate_wrapped_lines(title_block.get('text', ''), title_w, font_size)
        line_h = (font_size / 72.0) * 1.25
        title_text_h = max(0.22, lines * line_h + 0.06)
        title_bottom = title_top + title_text_h

        subtitle_top = subtitle_ph.top / EMU
        min_subtitle_top = title_bottom + MIN_GAP_IN
        if subtitle_top >= min_subtitle_top:
            return

        delta_in = min_subtitle_top - subtitle_top
        delta_emu = int(delta_in * EMU)
        subtitle_ph.top = subtitle_ph.top + delta_emu
    except Exception as e:
        print('_ensure_subtitle_placeholder_clears_title error:', e)


# Backward-compatible alias: build_slide historically called this helper by the
# older "title_overflow" name. Keep both spellings so template content slides
# cannot fail before render_blocks() runs.
def _shift_blocks_for_title_overflow(slide, blocks, use_template):
    return _shift_blocks_for_title_gap(slide, blocks, use_template)


# ─── BLOCKS-BASED RENDERER ────────────────────────────────────────────────────
# These thin wrappers are called by render_blocks() — the new pure renderer that
# iterates slide_spec['blocks'] and dispatches each to its typed handler.
# All layout decisions are pre-computed by flattenToBlocks() in agent5.js.

def render_block_title(slide, block, bt, use_template):
    """Render a title or subtitle block.

    Template mode — seed text into the master's placeholder so that font, size,
    colour, and position stay consistent across all content slides.  The caller
    must NOT have removed the title/subtitle placeholder before this point.

    Scratch mode — fall back to a free-form text box at the coordinates Agent 5
    computed.
    """
    text = block.get('text', '')
    if not text:
        return
    btype = block.get('block_type', 'title')
    if use_template:
        # Use the layout's title (idx 0) or subtitle (idx 1) placeholder.
        # preserve_template_style=True keeps the master's font/size/colour;
        # we only inject the text so every content slide looks identical.
        ph_idx = 0 if btype == 'title' else 1
        place_in_placeholder(slide, ph_idx, text, block, bt,
                             preserve_template_style=True,
                             compact_title=False)
        return
    # Scratch mode — render as a positioned text box
    add_text_box(slide,
        block.get('x', 0.4), block.get('y', 0.15),
        block.get('w', 9.2),  block.get('h', 0.7),
        text,
        block.get('font_family', bt.get('title_font_family', 'Arial')),
        block.get('font_size', 20),
        block.get('bold', True),
        block.get('color', bt.get('title_color', '#1A3C8F')),
        block.get('align', 'left'),
        block.get('valign', 'top'))


def render_block_text_box(slide, block, bt):
    """Render a generic text_box block."""
    text = block.get('text', '')
    if not text:
        return
    add_text_box(slide,
        block.get('x', 0), block.get('y', 0),
        block.get('w', 1),  block.get('h', 0.3),
        text,
        block.get('font_family', bt.get('body_font_family', 'Arial')),
        block.get('font_size', 10),
        block.get('bold', False),
        block.get('color', '#000000'),
        block.get('align', 'left'),
        block.get('valign', 'middle'))


def render_block_rect(slide, block):
    """Render a rect block (filled rectangle with optional border)."""
    fill   = block.get('fill_color')
    border = block.get('border_color')
    bw     = block.get('border_width', 0)
    cr     = block.get('corner_radius', 0)
    if not fill and not (border and bw):
        return
    add_filled_rect(slide,
        block.get('x', 0), block.get('y', 0),
        block.get('w', 1), block.get('h', 0.3),
        fill_hex=fill, border_hex=border,
        border_pt=bw, corner_radius=cr)


def render_block_circle(slide, block, bt):
    """Render a filled circle (oval) with centered bold text — used for circle_badge group headers."""
    x = block.get('x', 0)
    y = block.get('y', 0)
    w = block.get('w', 0.3)
    h = block.get('h', 0.3)
    fill = block.get('fill_color')
    shape = slide.shapes.add_shape(
        9,  # MSO_SHAPE.OVAL
        inches(x), inches(y), inches(w), inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill)
    else:
        shape.fill.background()
    shape.line.fill.background()

    text = str(block.get('text', ''))
    if text:
        tf = shape.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.bold = True
        run.font.size = Pt(block.get('font_size', 10))
        run.font.color.rgb = hex_to_rgb(block.get('font_color') or bt.get('body_color', '#FFFFFF'))
        fam = block.get('font_family') or bt.get('body_font_family', 'Arial')
        if fam:
            run.font.name = fam


def render_block_rule(slide, block, bt):
    """Render a horizontal rule line — uses a connector for true hairline weight."""
    from pptx.enum.shapes import MSO_CONNECTOR
    x = block.get('x', 0)
    y = block.get('y', 0)
    w = block.get('w', 1)
    color = block.get('color') or bt.get('primary_color', '#1A3C8F')
    line_width_pt = float(block.get('line_width', 0.5) or 0.5)
    try:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            inches(x), inches(y), inches(x + w), inches(y)
        )
        line.line.color.rgb = hex_to_rgb(color)
        line.line.width = pt(line_width_pt)
    except Exception:
        pass


def render_block_line(slide, block, bt):
    """Render a straight connector line between two endpoints."""
    from pptx.enum.shapes import MSO_CONNECTOR
    x1 = block.get('x1', block.get('x', 0))
    y1 = block.get('y1', block.get('y', 0))
    x2 = block.get('x2', x1)
    y2 = block.get('y2', y1)
    color = block.get('color') or bt.get('primary_color', '#1A3C8F')
    width_pt = float(block.get('line_width', 0.5) or 0.5)
    try:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            inches(x1), inches(y1), inches(x2), inches(y2)
        )
        line.line.color.rgb = hex_to_rgb(color)
        line.line.width = pt(width_pt)
    except Exception:
        pass


def render_block_bullet_list(slide, block, bt):
    """Render a bullet_list block — calls existing render_insight_text logic."""
    # Build a minimal artifact dict that render_insight_text expects.
    # suppress_heading=True because the header band was already emitted as
    # separate text_box/rect blocks before this bullet_list block.
    artifact = {
        'type':             'insight_text',
        'insight_mode':     'standard',
        'x':                block.get('x', 0),
        'y':                block.get('y', 0),
        'w':                block.get('w', 4),
        'h':                block.get('h', 3),
        'points':           block.get('points', []),
        'body_style':       dict(block.get('body_style', {})),
        'sentiment':        block.get('sentiment', 'neutral'),
        'heading':          None,
        'insight_header':   None,
        'style':            {},   # container rect already emitted as a separate rect block
    }
    # Apply padding if provided (for grouped bullet boxes)
    padding = block.get('padding', {})
    if padding:
        bs = artifact['body_style']
        if padding.get('top') is not None and not bs.get('padding_top'):
            bs['padding_top'] = padding.get('top')
        if padding.get('bottom') is not None and not bs.get('padding_bottom'):
            bs['padding_bottom'] = padding.get('bottom')
        if padding.get('left') is not None and not bs.get('padding_left'):
            bs['padding_left'] = padding.get('left')
        if padding.get('right') is not None and not bs.get('padding_right'):
            bs['padding_right'] = padding.get('right')
    elif renderer_fallback_allowed(block):
        bs = artifact['body_style']
        if not bs.get('padding_top'):    bs['padding_top']    = 0.08
        if not bs.get('padding_bottom'): bs['padding_bottom'] = 0.08
        if not bs.get('padding_left'):   bs['padding_left']   = 0.10
        if not bs.get('padding_right'):  bs['padding_right']  = 0.10
    artifact = apply_block_render_metadata(artifact, block, default_type='insight_text', default_subtype='standard')
    render_insight_text(slide, artifact, bt, suppress_heading=True)


def render_block_chart(slide, block, bt):
    """Render a chart block."""
    # Build artifact dict compatible with render_chart
    artifact = dict(block)
    artifact['type'] = 'chart'
    if not artifact.get('chart_type') and block.get('artifact_subtype'):
        artifact['chart_type'] = block.get('artifact_subtype')
    if not artifact.get('chart_header') and block.get('artifact_header_text'):
        artifact['chart_header'] = block.get('artifact_header_text')
    artifact = apply_block_render_metadata(artifact, block, default_type='chart', default_subtype=artifact.get('chart_type'))
    # Map pre-computed fields into _computed sub-dict for backward compat
    artifact['_computed'] = {
        'legend_position':         block.get('legend_position', 'none'),
        'data_label_size':         block.get('data_label_size', 9),
        'category_label_rotation': block.get('category_label_rotation', 0),
    }
    # Merge brand_tokens if provided in block
    block_bt = block.get('brand_tokens', {})
    merged_bt = {**bt, **{k: v for k, v in (block_bt or {}).items() if v}}
    render_chart(slide, artifact, merged_bt)


def render_block_table(slide, block, bt):
    """Render a table block."""
    artifact = dict(block)
    artifact['type'] = 'table'
    if not artifact.get('table_header') and block.get('artifact_header_text'):
        artifact['table_header'] = block.get('artifact_header_text')
    artifact = apply_block_render_metadata(artifact, block, default_type='table', default_subtype=block.get('artifact_subtype') or 'standard')
    render_table(slide, artifact, bt)


def render_block_workflow(slide, block, bt):
    """Render a workflow block."""
    artifact = dict(block)
    artifact['type'] = 'workflow'
    if not artifact.get('workflow_header') and block.get('artifact_header_text'):
        artifact['workflow_header'] = block.get('artifact_header_text')
    artifact = apply_block_render_metadata(artifact, block, default_type='workflow', default_subtype=block.get('artifact_subtype') or artifact.get('workflow_type'))
    render_workflow(slide, artifact, bt)


def render_block_footer(slide, block, bt):
    """Render footer / page_number text block."""
    text = block.get('text', '')
    if not text:
        return
    add_text_box(slide,
        block.get('x', 0.4), block.get('y', 7.3),
        block.get('w', 5.0),  block.get('h', 0.22),
        text,
        block.get('font_family', bt.get('body_font_family', 'Arial')),
        block.get('font_size', 8),
        False,
        block.get('color', '#AAAAAA'),
        block.get('align', 'left'),
        'middle')


def render_blocks(slide, slide_spec, bt, use_template):
    """
    Pure renderer: iterates slide_spec['blocks'] and dispatches each block
    to its type-specific render function. No layout decisions are made here.
    """
    blocks = list(slide_spec.get('blocks') or [])

    # Template mode needs a two-pass header render:
    # 1. place title
    # 2. ensure subtitle placeholder clears the wrapped title text
    # 3. place subtitle
    # 4. shift remaining content blocks based on the actual header text bottom
    if use_template and blocks:
        title_block = next((b for b in blocks if b.get('block_type') == 'title'), None)
        subtitle_block = next((b for b in blocks if b.get('block_type') == 'subtitle'), None)
        if title_block:
            try:
                render_block_title(slide, title_block, bt, use_template)
            except Exception as e:
                print(f'render_blocks: error on block_type=title: {e}')
        if subtitle_block:
            _ensure_subtitle_placeholder_clears_title(slide, blocks)
            try:
                render_block_title(slide, subtitle_block, bt, use_template)
            except Exception as e:
                print(f'render_blocks: error on block_type=subtitle: {e}')
        blocks = _shift_blocks_for_title_gap(slide, blocks, use_template)

    for block in blocks:
        btype = block.get('block_type', '')
        try:
            if use_template and btype in ('title', 'subtitle'):
                continue
            if btype in ('title', 'subtitle'):
                render_block_title(slide, block, bt, use_template)
            elif btype == 'text_box':
                render_block_text_box(slide, block, bt)
            elif btype == 'rect':
                render_block_rect(slide, block)
            elif btype == 'circle':
                render_block_circle(slide, block, bt)
            elif btype == 'rule':
                render_block_rule(slide, block, bt)
            elif btype == 'line':
                render_block_line(slide, block, bt)
            elif btype == 'bullet_list':
                render_block_bullet_list(slide, block, bt)
            elif btype == 'chart':
                render_block_chart(slide, block, bt)
            elif btype == 'table':
                render_block_table(slide, block, bt)
            elif btype == 'workflow':
                render_block_workflow(slide, block, bt)
            elif btype in ('footer', 'page_number'):
                render_block_footer(slide, block, bt)
            # 'image' (logo) blocks are not emitted by flattenToBlocks():
            # template master carries the logo; scratch mode renders it from global_elements.
        except Exception as e:
            print(f'render_blocks: error on block_type={btype}: {e}')


def _legacy_render_zones(slide, slide_spec, bt, use_template, layout_mode,
                          _ph_bounds, _content_ph_frames, slide_header_style,
                          title_shrink_in, cvs):
    """
    Legacy zones-based rendering path.
    Called from build_slide only when slide_spec has no 'blocks' key (old specs).
    """
    _zones = slide_spec.get('zones', [])
    if title_shrink_in > 0.02:
        for _z in _zones:
            for _a in (_z.get('artifacts') or []):
                if _a.get('y') is not None:
                    _a['y'] = round(max(0.0, _a['y'] - title_shrink_in), 4)
                if _a.get('header_block') and _a['header_block'].get('y') is not None:
                    _a['header_block']['y'] = round(
                        max(0.0, _a['header_block']['y'] - title_shrink_in), 4)

    _layout_ph_bounds = _ph_bounds if (layout_mode and use_template) else {}
    for zone_idx, zone in enumerate(_zones):
        zone = normalize_zone_artifact_stack(zone)
        _zone_arts_check = zone.get('artifacts', [])
        if (layout_mode
                and len(_zone_arts_check) >= 2
                and all(a.get('x') is None for a in _zone_arts_check)
                and _content_ph_frames):
            _slot = min(zone_idx, len(_content_ph_frames) - 1)
            _phf  = _content_ph_frames[_slot]
            _synth_zone = {
                **zone,
                'frame': {
                    'x': _phf['x'], 'y': _phf['y'],
                    'w': _phf['w'], 'h': _phf['h'],
                    'padding': {'top': 0.0, 'right': 0.0, 'bottom': 0.0, 'left': 0.0}
                }
            }
            zone = normalize_zone_artifact_stack(_synth_zone)
        hdr_ph_idx = zone.get('header_ph_idx') if layout_mode else None
        zone_artifacts = zone.get('artifacts', [])
        for art_idx, artifact in enumerate(zone_artifacts):
            ph_frame = None
            if layout_mode:
                ph_idx_spec = artifact.get('placeholder_idx')
                _wf_has_container = (
                    str(artifact.get('type', '')).lower() == 'workflow'
                    and isinstance(artifact.get('container'), dict)
                    and all(artifact.get('container', {}).get(k) is not None for k in ('x', 'y', 'w', 'h'))
                )
                needs_bounds = (
                    ph_idx_spec is not None
                    and len(zone_artifacts) == 1
                    and artifact.get('x') is None
                    and not _wf_has_container
                )
                if needs_bounds and ph_idx_spec in _layout_ph_bounds:
                    ph_frame = _layout_ph_bounds[ph_idx_spec]
                if ph_frame is None and artifact.get('x') is None and not _wf_has_container and _content_ph_frames:
                    fallback_slot = zone_idx if len(zone_artifacts) == 1 else (zone_idx + art_idx)
                    fallback_slot = min(fallback_slot, len(_content_ph_frames) - 1)
                    ph_frame = _content_ph_frames[fallback_slot]

            if ph_frame is None and artifact.get('x') is None and not (
                str(artifact.get('type', '')).lower() == 'workflow'
                and isinstance(artifact.get('container'), dict)
                and all(artifact.get('container', {}).get(k) is not None for k in ('x', 'y', 'w', 'h'))
            ):
                _zf = zone.get('frame') or {}
                if all(_zf.get(k) is not None for k in ('x', 'y', 'w', 'h')):
                    _zp = _zf.get('padding') or {}
                    _zl = float(_zp.get('left', 0) or 0)
                    _zt = float(_zp.get('top',  0) or 0)
                    _zr = float(_zp.get('right', 0) or 0)
                    _zb = float(_zp.get('bottom', 0) or 0)
                    ph_frame = {
                        'x': float(_zf['x']) + _zl,
                        'y': float(_zf['y']) + _zt,
                        'w': max(0.2, float(_zf['w']) - _zl - _zr),
                        'h': max(0.2, float(_zf['h']) - _zt - _zb),
                    }

            render_artifact(slide, artifact, bt, ph_frame=ph_frame, header_ph_idx=hdr_ph_idx,
                            header_style=slide_header_style,
                            slide_w=float(cvs.get('width_in') or 13.33),
                            slide_h=float(cvs.get('height_in') or 7.5))

    # Global elements (scratch only)
    if not use_template:
        ge = slide_spec.get('global_elements', {})
        logo = ge.get('logo', {})
        if logo.get('show') and logo.get('image_base64'):
            add_image_box(slide, logo['image_base64'],
                logo.get('x', 0.0), logo.get('y', 0.0),
                logo.get('w', 1.2), logo.get('h', 0.4))
        footer = ge.get('footer', {})
        if footer.get('show'):
            add_text_box(slide,
                footer.get('x', 0.4), footer.get('y', 7.5),
                footer.get('w', 5),   footer.get('h', 0.25),
                'Confidential',
                footer.get('font_family', 'Arial'),
                footer.get('font_size', 8),
                False, footer.get('color', '#AAAAAA'),
                footer.get('align', 'left'), 'middle')
        pn = ge.get('page_number', {})
        if pn.get('show'):
            add_text_box(slide,
                pn.get('x', 9.5), pn.get('y', 7.5),
                pn.get('w', 0.8), pn.get('h', 0.25),
                str(slide_spec.get('slide_number', '')),
                pn.get('font_family', 'Arial'),
                pn.get('font_size', 8),
                False, pn.get('color', '#AAAAAA'),
                'right', 'middle')


def build_slide(prs, slide_spec, blank_layout, use_template=False,
                title_layout_name=None, divider_layout_name=None):
    """
    Build one slide from its spec object.

    Three rendering paths:
    1. Title / divider  — template mode: text into placeholders only.
    2. Content          — template mode: title/subtitle into placeholders (idx 0/1);
                          artifacts rendered at coordinates OR at placeholder bounds
                          when layout_mode=True and placeholder_idx is set.
    3. Any type         — scratch mode (use_template=False): full coordinates for
                          everything including title, subtitle, logo, footer, page num.
    """
    slide_type           = slide_spec.get('slide_type', 'content')
    layout_mode          = slide_spec.get('layout_mode', False)
    selected_layout_name = slide_spec.get('selected_layout_name', '')
    bt                   = slide_spec.get('brand_tokens', {})
    cvs                  = slide_spec.get('canvas', {})
    tb                   = dict(slide_spec.get('title_block') or {})
    sb                   = dict(slide_spec.get('subtitle_block') or {})
    if not tb.get('text') and slide_spec.get('title'):
        tb['text'] = slide_spec.get('title')
    if not sb.get('text') and slide_spec.get('subtitle'):
        sb['text'] = slide_spec.get('subtitle')
    slide_header_style   = infer_slide_header_style(slide_spec)

    # ── Choose layout ────────────────────────────────────────────────────────
    layout = blank_layout
    if use_template:
        if selected_layout_name and slide_type == 'content':
            # Named layout selected by Agent 4 — only valid for content slides.
            # Title/divider/thank-you slides must never use a content-area layout;
            # they fall through to their dedicated type handlers below.
            named = find_layout_by_name(prs, selected_layout_name)
            if named:
                layout = named
                print(f'  Slide {slide_spec.get("slide_number","?")}: layout="{named.name}"')
            else:
                # Named layout not found — fall back to neutral content layout
                layout = find_content_fallback_layout(prs) or blank_layout
                print(f'  Slide {slide_spec.get("slide_number","?")}: layout "{selected_layout_name}" not found — using content fallback')
        elif slide_type == 'content':
            # Scratch-mode content slide (layout_mode=False or no selected_layout_name):
            # use a neutral content layout, never a title/divider/closing layout.
            layout = find_content_fallback_layout(prs) or blank_layout

        elif slide_type == 'title':
            # Use the explicit name from Agent 2's rulebook first (most reliable),
            # then fall back to OOXML type="title" matching, then name heuristics.
            if title_layout_name:
                layout = find_layout_by_name(prs, title_layout_name) or blank_layout
            else:
                for lyt in prs.slide_layouts:
                    # python-pptx exposes the OOXML type via the internal XML
                    try:
                        import pptx.oxml.ns as _ns
                        lyt_type = lyt._element.get('type', '')
                    except Exception:
                        lyt_type = ''
                    if lyt_type == 'title':
                        layout = lyt; break
                else:
                    for lyt in prs.slide_layouts:
                        lname = (lyt.name or '').lower()
                        if 'title' in lname and 'section' not in lname and 'content' not in lname:
                            layout = lyt; break

        elif slide_type == 'divider':
            # Same priority: explicit rulebook name → OOXML secHead → name heuristics.
            if divider_layout_name:
                layout = find_layout_by_name(prs, divider_layout_name) or blank_layout
            else:
                for lyt in prs.slide_layouts:
                    try:
                        import pptx.oxml.ns as _ns
                        lyt_type = lyt._element.get('type', '')
                    except Exception:
                        lyt_type = ''
                    if lyt_type == 'secHead':
                        layout = lyt; break
                else:
                    for lyt in prs.slide_layouts:
                        lname = (lyt.name or '').lower()
                        if 'section' in lname or 'divider' in lname:
                            layout = lyt; break

    slide = prs.slides.add_slide(layout)
    # Always initialise these so layout_mode references below are safe
    # even when use_template=False (scratch mode).
    _ph_bounds         = {}
    _content_ph_frames = []
    if use_template:
        from pptx.enum.shapes import PP_PLACEHOLDER as _PH
        _sanitise_system_placeholders(slide, slide_spec.get('slide_number', ''))
        # ── Snapshot placeholder bounds BEFORE any removal ───────────────────
        # render_artifact needs the bounds of content placeholders (12, 13, …) for
        # positioning, but those placeholders must be removed to prevent ghost text.
        # We read all bounds now, then remove, then pass the snapshot to renderers.
        _ph_bounds = {}   # idx → {'x','y','w','h'} in inches
        _content_ph_frames = []  # ordered fallback slots for layout_mode when placeholder_idx is missing
        _SYSTEM_TYPES = {_PH.TITLE, _PH.CENTER_TITLE, _PH.SUBTITLE,
                         _PH.DATE, _PH.FOOTER, _PH.SLIDE_NUMBER}
        for _ph in slide.placeholders:
            try:
                _ph_idx = _ph.placeholder_format.idx
                _ph_type = _ph.placeholder_format.type
                _bounds = {
                    'x': _ph.left   / 914400,
                    'y': _ph.top    / 914400,
                    'w': _ph.width  / 914400,
                    'h': _ph.height / 914400,
                }
                _ph_bounds[_ph_idx] = _bounds
                # Some Agent 5 / 5.1 outputs reach Agent 6 without placeholder_idx.
                # Keep an ordered list of likely content slots so layout-mode content
                # can still render into the correct body area instead of disappearing.
                if (_ph_type not in _SYSTEM_TYPES and
                    _bounds['w'] >= 1.0 and _bounds['h'] >= 0.5):
                    _content_ph_frames.append({
                        'idx': _ph_idx,
                        **_bounds
                    })
            except Exception:
                pass
        _content_ph_frames.sort(key=lambda p: (round(p['y'] * 2), p['x'], p['idx']))

        # Collect header placeholder indices we will explicitly write heading text to.
        # These must survive removal so _write_heading_to_header_ph can use them.
        _header_ph_keep = set()
        for _z in slide_spec.get('zones', []):
            _hpi = _z.get('header_ph_idx')
            if _hpi is not None:
                _header_ph_keep.add(int(_hpi))
        _remove_content_placeholders(slide, keep_indices=_header_ph_keep)

    # ── Background (scratch only — master handles it in template mode) ────────
    if not use_template:
        bg_color = (cvs.get('background') or {}).get('color', '#FFFFFF')
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex_to_rgb(bg_color)

    # ════════════════════════════════════════════════════════════════════════
    # PATH 1 — TITLE & DIVIDER
    # Only write text content; the layout provides all visual decoration.
    # ════════════════════════════════════════════════════════════════════════
    _has_blocks = bool(slide_spec.get('blocks'))
    if slide_type in ('title', 'divider') and _has_blocks and not use_template:
        render_blocks(slide, slide_spec, bt, use_template)
        if use_template:
            _remove_empty_placeholders(slide)
        _write_speaker_note(slide, slide_spec.get('speaker_note', ''))
        return slide
    if slide_type in ('title', 'divider'):
        if use_template:
            if tb.get('text'):
                place_in_placeholder(
                    slide, 0, tb['text'], tb, bt,
                    preserve_template_style=True,
                    compact_title=False
                )
            if sb.get('text'):
                place_in_placeholder(
                    slide, 1, sb['text'], sb, bt,
                    preserve_template_style=True,
                    compact_title=False
                )
        else:
            # Scratch mode — free-form text boxes
            if tb.get('text'):
                add_text_box(slide,
                    tb.get('x', 0.4), tb.get('y', 0.15),
                    tb.get('w', 9.0), tb.get('h', 0.8),
                    tb['text'],
                    tb.get('font_family', bt.get('title_font_family', 'Arial')),
                    tb.get('font_size', 28),
                    tb.get('font_weight', 'bold') in ('bold', 'semibold'),
                    tb.get('color', bt.get('title_color', '#1A3C8F')),
                    tb.get('align', 'center'), 'middle')
            if sb.get('text'):
                add_text_box(slide,
                    sb.get('x', 0.4), sb.get('y', 1.2),
                    sb.get('w', 9.0), sb.get('h', 0.6),
                    sb['text'],
                    sb.get('font_family', bt.get('body_font_family', 'Arial')),
                    sb.get('font_size', 16),
                    False,
                    sb.get('color', bt.get('body_color', '#333333')),
                    sb.get('align', 'center'), 'middle')
        if use_template:
            _remove_empty_placeholders(slide)
        _write_speaker_note(slide, slide_spec.get('speaker_note', ''))
        return slide

    # ════════════════════════════════════════════════════════════════════════
    # PATH 2 & 3 — CONTENT SLIDES
    # ════════════════════════════════════════════════════════════════════════

    # Title & subtitle: rendered here only in the legacy (no-blocks) path.
    # When blocks[] is present, render_blocks dispatches title/subtitle blocks
    # via render_block_title — rendering them here too would double-write.
    # title_shrink_in: how much the title placeholder was compacted (inches).
    # Legacy zone artifacts are positioned assuming the full placeholder height,
    # so we shift their y-coords up by this amount to close the resulting gap.
    # (Not applied when using blocks[] — blocks have absolute pre-computed coords.)
    title_shrink_in = 0.0
    if not _has_blocks:
        if use_template:
            if tb.get('text'):
                title_shrink_in = place_in_placeholder(slide, 0, tb['text'], tb, bt) or 0.0
            if sb.get('text'):
                place_in_placeholder(slide, 1, sb['text'], sb, bt)
        else:
            if tb.get('text'):
                add_text_box(slide,
                    tb.get('x', 0.4), tb.get('y', 0.15),
                    tb.get('w', 9.0), tb.get('h', 0.8),
                    tb['text'],
                    tb.get('font_family', bt.get('title_font_family', 'Arial')),
                    tb.get('font_size', 20),
                    tb.get('font_weight', 'bold') in ('bold', 'semibold'),
                    tb.get('color', bt.get('title_color', '#1A3C8F')),
                    tb.get('align', 'left'), 'top')
                # Scratch mode: shrink title box to actual text height so blank space
                # below the title doesn't inflate the visual gap to the content zone.
                try:
                    _tf = slide.shapes[-1].text_frame
                    _tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    _tf.vertical_anchor = MSO_ANCHOR.TOP
                except Exception:
                    pass
            if sb.get('text'):
                add_text_box(slide,
                    sb.get('x', 0.4), sb.get('y', 1.0),
                    sb.get('w', 9.0), sb.get('h', 0.5),
                    sb['text'],
                    sb.get('font_family', bt.get('body_font_family', 'Arial')),
                    sb.get('font_size', 14),
                    sb.get('font_weight', 'semibold') in ('bold', 'semibold'),
                    sb.get('color', bt.get('body_color', '#333333')),
                    sb.get('align', 'left'), 'middle')

    # ════════════════════════════════════════════════════════════════════════
    # BLOCKS PATH ONLY
    # Agent 5 is responsible for final geometry via flattenToBlocks().
    # Agent 6 / generate_pptx.py must act as a pure renderer for content slides.
    # ════════════════════════════════════════════════════════════════════════
    if not _has_blocks:
        raise ValueError(
            f"Content slide {slide_spec.get('slide_number', '?')} is missing finalized blocks[]. "
            "Agent 5 must provide fully flattened render blocks for Agent 6."
        )

    # Title, subtitle, artifacts, and global_elements are all pre-flattened
    # into blocks[] by flattenToBlocks() in agent5.js.
    # Template mode header clearance is handled inside render_blocks() after the
    # title/subtitle placeholders are actually populated.
    render_blocks(slide, slide_spec, bt, use_template)

    if use_template:
        _remove_empty_placeholders(slide)
    _write_speaker_note(slide, slide_spec.get('speaker_note', ''))
    return slide


# ─── MAIN BUILDER ─────────────────────────────────────────────────────────────

def build_presentation(final_spec, brand_rulebook, template_b64=None):
    """
    Build a complete Presentation from the Agent 5 final spec.

    When template_b64 is provided the function:
      1. Loads the brand PPTX template (preserving slide masters)
      2. Strips all existing content slides from it
      3. Adds new slides backed by the template's blank layout

    This keeps the full brand identity (master background, decorative shapes,
    registered fonts, theme colours) without the need for Agent 5 to manually
    specify every design token.
    """
    if not final_spec:
        raise ValueError('finalSpec is empty')
    invalid_slides = []
    for slide_spec in final_spec:
        slide_no = slide_spec.get('slide_number', '?')
        has_canvas = bool(slide_spec.get('canvas'))
        has_blocks = isinstance(slide_spec.get('blocks'), list) and len(slide_spec.get('blocks', [])) > 0
        if not has_canvas or not has_blocks:
            invalid_slides.append(str(slide_no))
    if invalid_slides:
        raise ValueError(
            'Agent 6 requires Agent 5 render specs with canvas + non-empty blocks[] on every slide. '
            'Invalid slides: ' + ', '.join(invalid_slides)
        )

    use_template = bool(template_b64)

    if use_template:
        print('build_presentation: loading brand template…')
        try:
            prs = load_template_prs(template_b64)
            print(f'  Template loaded — {len(prs.slide_masters)} master(s), {len(prs.slide_layouts)} layout(s)')
        except Exception as e:
            print('  Template load failed, falling back to blank presentation:', e)
            prs = Presentation()
            use_template = False
    else:
        prs = Presentation()

    # Slide dimensions — use template's native size if template is loaded,
    # otherwise take from the first slide's canvas spec
    if not use_template:
        first_canvas = (final_spec[0].get('canvas') or {})
        w_in = float(first_canvas.get('width_in')  or 11.02)
        h_in = float(first_canvas.get('height_in') or 8.29)
        prs.slide_width  = Inches(w_in)
        prs.slide_height = Inches(h_in)

    blank_layout = find_blank_layout(prs)
    content_fallback_layout = find_content_fallback_layout(prs)
    print(f'  Using fallback layout: "{blank_layout.name}" (use_template={use_template})')
    if content_fallback_layout is not None:
        print(f'  Content fallback layout: "{content_fallback_layout.name}"')

    # Extract explicit layout names from Agent 2 rulebook so build_slide can
    # select title / divider layouts reliably without keyword heuristics.
    title_layout_name   = brand_rulebook.get('title_layout_name')
    divider_layout_name = brand_rulebook.get('divider_layout_name')
    if title_layout_name:
        print(f'  Title layout:   "{title_layout_name}"')
    if divider_layout_name:
        print(f'  Divider layout: "{divider_layout_name}"')

    for slide_spec in final_spec:
        try:
            build_slide(prs, slide_spec, blank_layout, use_template,
                        title_layout_name=title_layout_name,
                        divider_layout_name=divider_layout_name)
        except Exception as e:
            print(f"Error building slide {slide_spec.get('slide_number', '?')}:", e)
            traceback.print_exc()
            # Do NOT add a blank slide here — prs.slides.add_slide already ran
            # inside build_slide before the exception, so the slide exists.
            # Adding another blank here is what caused the slide count to balloon.

    return prs


# ─── VERCEL HANDLER ──────────────────────────────────────────────────────────

class handler(BaseHTTPRequestHandler):

    def do_POST(self):
        try:
            length = int(self.headers.get('Content-Length', 0))
            body   = self.rfile.read(length)
            data   = json.loads(body)

            final_spec     = data.get('finalSpec', [])
            brand_rulebook = data.get('brandRulebook', {})
            template_b64   = data.get('templateB64') or None

            if not final_spec:
                self._json(400, {'success': False, 'error': 'finalSpec is missing or empty'})
                return

            # Build the presentation (optionally from a brand template)
            prs = build_presentation(final_spec, brand_rulebook, template_b64)

            # Serialize to bytes
            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)
            pptx_bytes = buf.read()

            # Encode as base64
            b64 = base64.b64encode(pptx_bytes).decode('utf-8')

            # Filename
            date_str = datetime.now().strftime('%Y%m%d')
            title_raw = (final_spec[0].get('title') or 'presentation')
            title_slug = ''.join(c if c.isalnum() else '_' for c in title_raw.lower())[:30]
            filename = title_slug + '_' + date_str + '.pptx'

            self._json(200, {
                'success':  True,
                'data':     b64,
                'slides':   len(final_spec),
                'filename': filename
            })

        except Exception as e:
            traceback.print_exc()
            self._json(500, {'success': False, 'error': str(e)})

    def _json(self, status, payload):
        body = json.dumps(payload).encode('utf-8')
        self.send_response(status)
        self.send_header('Content-Type',   'application/json')
        self.send_header('Content-Length', str(len(body)))
        self.send_header('Access-Control-Allow-Origin', '*')
        self.end_headers()
        self.wfile.write(body)

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin',  '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()

    def log_message(self, format, *args):
        pass  # suppress default request logging
